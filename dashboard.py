import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import os
from datetime import datetime, date, timedelta
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

st.set_page_config(
    page_title="门店运营看板",
    page_icon="🏪",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── 样式（浅色专业风）─────────────────────────────────────────────────────────
st.markdown("""
<style>
[data-testid="stAppViewContainer"] { background: #f4f6fa; }
[data-testid="stSidebar"] { background: #ffffff; border-right: 1px solid #e2e8f0; }
[data-testid="stSidebar"] * { color: #1a202c !important; }
[data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2,
[data-testid="stSidebar"] h3, [data-testid="stSidebar"] h4 { color: #1a202c !important; }

.block-container { padding-top: 1.5rem !important; }

.metric-card {
    background: #ffffff;
    border: 1px solid #e2e8f0;
    border-radius: 10px;
    padding: 16px 20px;
    margin-bottom: 10px;
    box-shadow: 0 1px 3px rgba(0,0,0,0.04);
}
.metric-label  { color: #64748b; font-size: 13px; margin-bottom: 6px; font-weight: 500; }
.metric-value  { color: #0f172a; font-size: 26px; font-weight: 700; line-height: 1.2; }
.metric-sub    { color: #94a3b8; font-size: 11px; margin-top: 4px; }

.section-title {
    color: #1e293b;
    font-size: 16px;
    font-weight: 700;
    margin: 24px 0 12px;
    padding-left: 10px;
    border-left: 4px solid #4f46e5;
}

div[data-testid="stTabs"] button {
    color: #64748b;
    font-size: 15px;
    font-weight: 500;
}
div[data-testid="stTabs"] button[aria-selected="true"] {
    color: #4f46e5;
    border-bottom: 2px solid #4f46e5;
}

/* 主区域文字 */
h1, h2, h3, h4, p, label, span { color: #1e293b; }

/* 顶部信息条 */
.top-info {
    background: #ffffff;
    border: 1px solid #e2e8f0;
    border-radius: 10px;
    padding: 12px 20px;
    margin-bottom: 16px;
}
.top-info b { color: #4f46e5; }

/* 日报头部 */
.report-header {
    background: linear-gradient(135deg, #4f46e5, #6366f1);
    border-radius: 12px;
    padding: 20px 28px;
    margin-bottom: 16px;
    color: #ffffff;
}
.report-header h3 { color: #ffffff !important; margin: 0 0 8px; }
.report-header p, .report-header b { color: #ffffff !important; }
</style>
""", unsafe_allow_html=True)

PLATFORM_COLORS = {"美团": "#FFB300", "饿了么": "#1E88E5"}
CHART_BG = "#ffffff"
TEXT_COLOR = "#334155"
GRID_COLOR = "#e2e8f0"


def hex_to_rgba(hex_color: str, alpha: float = 0.25) -> str:
    """安全地把 #RRGGBB 转成 rgba(r,g,b,a)；任何异常都返回中性灰。"""
    try:
        h = str(hex_color).lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) >= 6:
            r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
            return f"rgba({r},{g},{b},{alpha})"
    except Exception:
        pass
    return f"rgba(100,116,139,{alpha})"


# ── 数据加载 ──────────────────────────────────────────────────────────────────
def load_meituan(file) -> pd.DataFrame:
    if isinstance(file, str):
        df = pd.read_csv(file, encoding="gbk")
    else:
        raw = file.read()
        df = pd.read_csv(io.BytesIO(raw), encoding="gbk")
    df["日期"] = pd.to_datetime(df["日期"], format="%Y%m%d")
    df["平台"] = "美团"
    df = df.rename(columns={
        "营业收入": "收入",
        "实付单均价": "单均实付",
        "入店人数": "进店人数",
        "入店转化率": "进店转化率",
        "入店新客": "新客进店人数",
        "入店老客": "老客进店人数",
        "入店次数": "进店次数",
        "平台服务费(含佣金和配送服务费)": "平台技术服务费",
        "顾客实付": "顾客实付总额",
    })
    # 统一字段
    df["商家补贴"] = pd.to_numeric(df.get("商家活动支出", 0), errors="coerce").fillna(0)
    df["平台费"]   = pd.to_numeric(df.get("平台技术服务费", 0), errors="coerce").fillna(0)
    df["优惠前总额"] = pd.to_numeric(df.get("优惠前总额", 0), errors="coerce").fillna(0)
    df["无效订单"] = np.nan        # 美团此表无该字段
    df["商户原因无效订单"] = np.nan
    df["退单费用"] = np.nan
    df["流量类型_当日"] = df.get("当日流量类型", pd.NA)
    df["流量类型_7日"]  = df.get("7日流量类型", pd.NA)
    df["流量类型_30日"] = df.get("30日流量类型", pd.NA)
    return df


def load_eleme(file) -> pd.DataFrame:
    if isinstance(file, str):
        df = pd.read_excel(file, sheet_name="data", header=0)
    else:
        df = pd.read_excel(file, sheet_name="data", header=0)
    df["日期"] = pd.to_datetime(df["日期"])
    df["平台"] = "饿了么"
    # 统一字段
    df["商家补贴"] = (pd.to_numeric(df.get("活动补贴", 0), errors="coerce").fillna(0)
                  + pd.to_numeric(df.get("代金券补贴", 0), errors="coerce").fillna(0))
    df["平台费"]   = (pd.to_numeric(df.get("平台技术服务费", 0), errors="coerce").fillna(0)
                  + pd.to_numeric(df.get("履约技术服务费", 0), errors="coerce").fillna(0))
    df["优惠前总额"] = pd.to_numeric(df.get("营业额", 0), errors="coerce").fillna(0)
    df["无效订单"]         = pd.to_numeric(df.get("无效订单", 0), errors="coerce").fillna(0)
    df["商户原因无效订单"] = pd.to_numeric(df.get("商户原因无效订单数", 0), errors="coerce").fillna(0)
    df["退单费用"]         = pd.to_numeric(df.get("退单费用", 0), errors="coerce").fillna(0)
    # 饿了么无平台流量类型字段
    df["流量类型_当日"] = pd.NA
    df["流量类型_7日"]  = pd.NA
    df["流量类型_30日"] = pd.NA
    return df


def load_all(mt_file, elm_file) -> pd.DataFrame:
    frames = []
    if mt_file is not None:
        try:
            frames.append(load_meituan(mt_file))
        except Exception as e:
            st.error(f"美团数据加载失败: {e}")
    if elm_file is not None:
        try:
            frames.append(load_eleme(elm_file))
        except Exception as e:
            st.error(f"饿了么数据加载失败: {e}")
    if not frames:
        return pd.DataFrame()

    shared_cols = [
        "日期", "门店名称", "平台", "收入", "有效订单",
        "单均实付", "顾客实付总额",
        "曝光人数", "进店人数", "下单转化率", "进店转化率",
        "曝光新客", "新客进店人数",
        "曝光老客", "老客进店人数",
        "曝光次数", "进店次数",
        "商家补贴", "平台费", "优惠前总额",
        "无效订单", "商户原因无效订单", "退单费用",
        "流量类型_当日", "流量类型_7日", "流量类型_30日",
    ]
    dfs = []
    for d in frames:
        available = [c for c in shared_cols if c in d.columns]
        dfs.append(d[available].copy())
    combined = pd.concat(dfs, ignore_index=True)
    numeric_cols = [
        "收入", "有效订单", "单均实付", "曝光人数", "进店人数",
        "下单转化率", "进店转化率", "曝光新客", "新客进店人数",
        "曝光老客", "老客进店人数", "曝光次数", "进店次数",
        "商家补贴", "平台费", "优惠前总额",
        "无效订单", "商户原因无效订单", "退单费用",
    ]
    for col in numeric_cols:
        if col in combined.columns:
            combined[col] = pd.to_numeric(combined[col], errors="coerce").fillna(0)
    return combined


# ── 工具函数 ──────────────────────────────────────────────────────────────────
def fmt_money(v):
    if v >= 10000:
        return f"¥{v/10000:.2f}万"
    return f"¥{v:,.0f}"

def fmt_int(v):
    return f"{int(v):,}"

def fmt_pct(v):
    return f"{v*100:.1f}%"


def kpi_card(label, value, sub=""):
    sub_html = f'<div class="metric-sub">{sub}</div>' if sub else ""
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-label">{label}</div>
        <div class="metric-value">{value}</div>
        {sub_html}
    </div>
    """, unsafe_allow_html=True)


def _base_layout(fig, height=None):
    fig.update_layout(
        plot_bgcolor=CHART_BG, paper_bgcolor=CHART_BG,
        font=dict(color=TEXT_COLOR, family="-apple-system, Microsoft YaHei, sans-serif", size=12),
        legend=dict(orientation="h", y=1.08, x=0, bgcolor="rgba(0,0,0,0)"),
        margin=dict(l=10, r=20, t=40, b=10),
        hovermode="x unified",
    )
    if height:
        fig.update_layout(height=height)
    return fig


# ── 图表 ──────────────────────────────────────────────────────────────────────
def plot_daily_trend(df):
    daily = df.groupby(["日期", "平台"]).agg(
        收入=("收入", "sum"),
        订单=("有效订单", "sum"),
    ).reset_index()
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    for plat in daily["平台"].unique():
        sub = daily[daily["平台"] == plat]
        color = PLATFORM_COLORS.get(plat, "#888")
        fig.add_trace(go.Scatter(
            x=sub["日期"], y=sub["收入"], name=f"{plat} 营业收入",
            line=dict(color=color, width=3),
            mode="lines+markers", marker=dict(size=8, line=dict(color="white", width=1.5)),
            hovertemplate=f"{plat} 收入: ¥%{{y:,.2f}}<extra></extra>",
        ), secondary_y=False)
        fig.add_trace(go.Bar(
            x=sub["日期"], y=sub["订单"], name=f"{plat} 订单量",
            marker_color=color, opacity=0.35,
            hovertemplate=f"{plat} 订单: %{{y:,.0f}}<extra></extra>",
        ), secondary_y=True)
    _base_layout(fig, height=420)
    fig.update_xaxes(gridcolor=GRID_COLOR, tickformat="%m-%d", showline=True, linecolor=GRID_COLOR)
    fig.update_yaxes(title_text="营业收入 (¥)", gridcolor=GRID_COLOR, tickformat=",.0f",
                     secondary_y=False, showline=True, linecolor=GRID_COLOR)
    fig.update_yaxes(title_text="订单量", gridcolor=GRID_COLOR, tickformat=",.0f",
                     secondary_y=True, showgrid=False)
    return fig


def _short_store_name(s):
    import re
    m = re.search(r"[（(](.+?)[)）]", s)
    return m.group(1) if m else s


def plot_store_rank(df, metric, label):
    store = df.groupby(["门店名称", "平台"]).agg(v=(metric, "sum")).reset_index()
    store = store.sort_values("v", ascending=True)
    store["门店短名"] = store["门店名称"].apply(_short_store_name)
    fig = px.bar(
        store, x="v", y="门店短名", color="平台",
        color_discrete_map=PLATFORM_COLORS,
        orientation="h", text="v",
        labels={"v": label, "门店短名": ""},
    )
    fig.update_traces(texttemplate="%{x:,.0f}", textposition="outside",
                      textfont=dict(color=TEXT_COLOR, size=11))
    _base_layout(fig, height=max(360, len(store) * 30))
    fig.update_xaxes(gridcolor=GRID_COLOR, tickformat=",.0f", showline=True, linecolor=GRID_COLOR)
    fig.update_yaxes(gridcolor=GRID_COLOR, automargin=True)
    fig.update_layout(bargap=0.35, margin=dict(l=10, r=80, t=40, b=10))
    return fig


def plot_funnel(df):
    exposure = df["曝光人数"].sum()
    visit    = df["进店人数"].sum()
    orders   = df["有效订单"].sum()
    fig = go.Figure(go.Funnel(
        y=["曝光人数", "进店人数", "下单人数"],
        x=[exposure, visit, orders],
        textinfo="value+percent initial",
        texttemplate="%{value:,.0f}<br>%{percentInitial}",
        textfont=dict(color="white", size=13),
        marker=dict(color=["#4f46e5", "#10b981", "#f59e0b"]),
        connector=dict(line=dict(color="#cbd5e1", width=2)),
    ))
    _base_layout(fig, height=380)
    fig.update_layout(margin=dict(l=10, r=10, t=20, b=10))
    return fig


def plot_new_old(df):
    new_exp = df["曝光新客"].sum() if "曝光新客" in df.columns else 0
    old_exp = df["曝光老客"].sum() if "曝光老客" in df.columns else 0
    new_vis = df["新客进店人数"].sum() if "新客进店人数" in df.columns else 0
    old_vis = df["老客进店人数"].sum() if "老客进店人数" in df.columns else 0

    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{"type": "pie"}, {"type": "pie"}]],
        subplot_titles=["曝光：新客 vs 老客", "进店：新客 vs 老客"],
    )
    for i, (n, o) in enumerate([(new_exp, old_exp), (new_vis, old_vis)], 1):
        fig.add_trace(go.Pie(
            labels=["新客", "老客"], values=[n, o],
            marker=dict(colors=["#4f46e5", "#10b981"], line=dict(color="white", width=2)),
            hole=0.55,
            textinfo="label+percent",
            texttemplate="%{label}<br>%{value:,.0f}<br>(%{percent})",
            textfont=dict(color="white", size=12),
        ), row=1, col=i)
    _base_layout(fig, height=380)
    fig.update_layout(showlegend=False)
    for ann in fig["layout"]["annotations"]:
        ann["font"] = dict(color=TEXT_COLOR, size=13)
    return fig


def plot_conversion_by_store(df):
    conv = df.groupby(["门店名称", "平台"]).agg(
        进店转化=("进店转化率", "mean"),
        下单转化=("下单转化率", "mean"),
    ).reset_index()
    conv["门店短名"] = conv["门店名称"].apply(_short_store_name)
    conv = conv.sort_values("下单转化", ascending=True)
    fig = go.Figure()
    fig.add_trace(go.Bar(
        name="进店转化率", y=conv["门店短名"], x=conv["进店转化"]*100,
        orientation="h", marker_color="#4f46e5", opacity=0.55,
        text=[f"{v*100:.1f}%" for v in conv["进店转化"]],
        textposition="outside", textfont=dict(color=TEXT_COLOR, size=10),
    ))
    fig.add_trace(go.Bar(
        name="下单转化率", y=conv["门店短名"], x=conv["下单转化"]*100,
        orientation="h", marker_color="#10b981",
        text=[f"{v*100:.1f}%" for v in conv["下单转化"]],
        textposition="outside", textfont=dict(color=TEXT_COLOR, size=10),
    ))
    _base_layout(fig, height=max(380, len(conv) * 32))
    fig.update_layout(barmode="group", margin=dict(l=10, r=80, t=40, b=10))
    fig.update_xaxes(gridcolor=GRID_COLOR, ticksuffix="%", showline=True, linecolor=GRID_COLOR)
    fig.update_yaxes(gridcolor=GRID_COLOR, automargin=True)
    return fig


# ── 预警计算 ──────────────────────────────────────────────────────────────────
ALERT_METRICS = {
    "营业收入":   {"col": "收入",        "fmt": lambda x: f"¥{x:,.2f}",        "low_is_bad": True},
    "有效订单":   {"col": "有效订单",     "fmt": lambda x: f"{int(x):,}",       "low_is_bad": True},
    "进店转化率": {"col": "进店转化率",   "fmt": lambda x: f"{x*100:.2f}%",     "low_is_bad": True},
    "下单转化率": {"col": "下单转化率",   "fmt": lambda x: f"{x*100:.2f}%",     "low_is_bad": True},
    "曝光人数":   {"col": "曝光人数",     "fmt": lambda x: f"{int(x):,}",       "low_is_bad": True},
    "进店人数":   {"col": "进店人数",     "fmt": lambda x: f"{int(x):,}",       "low_is_bad": True},
}


def compute_alerts(df, z_thresh=2.0, drop_thresh=0.30, monitor_metrics=None):
    """
    基于历史均值的异常检测：
      - Z-score：(今日值 - 历史均值) / 历史标准差，|z| > z_thresh 触发
      - 环比下跌：当指标"低=坏"时，下跌幅度 > drop_thresh 触发
    """
    if monitor_metrics is None:
        monitor_metrics = list(ALERT_METRICS.keys())

    alerts = []
    for (store, platform), grp in df.groupby(["门店名称", "平台"]):
        grp = grp.sort_values("日期").reset_index(drop=True)
        if len(grp) < 3:
            continue
        for i in range(len(grp)):
            today = grp.iloc[i]
            others = grp.drop(i)  # 除今日外所有历史数据
            if len(others) < 2:
                continue

            for name in monitor_metrics:
                info = ALERT_METRICS[name]
                col = info["col"]
                if col not in grp.columns:
                    continue
                today_val = today[col]
                hist_mean = others[col].mean()
                hist_std  = others[col].std()

                if hist_mean == 0 or pd.isna(hist_mean):
                    continue

                z = (today_val - hist_mean) / hist_std if hist_std and hist_std > 0 else 0
                deviation = (today_val - hist_mean) / hist_mean  # 相对偏离

                bad_direction = (deviation < 0) if info["low_is_bad"] else (deviation > 0)
                trigger_z    = abs(z) >= z_thresh and bad_direction
                trigger_drop = info["low_is_bad"] and (deviation <= -drop_thresh) and abs(deviation) > 0.05

                if trigger_z or trigger_drop:
                    if abs(z) >= 2.5 or abs(deviation) >= 0.5:
                        severity = "🔴 严重"
                        sev_rank = 3
                    elif abs(z) >= 2.0 or abs(deviation) >= 0.35:
                        severity = "🟠 警告"
                        sev_rank = 2
                    else:
                        severity = "🟡 注意"
                        sev_rank = 1

                    alerts.append({
                        "日期": today["日期"].date(),
                        "门店": store,
                        "平台": platform,
                        "指标": name,
                        "当日值": info["fmt"](today_val),
                        "历史均值": info["fmt"](hist_mean),
                        "偏离": f"{deviation*100:+.1f}%",
                        "Z值": f"{z:+.2f}",
                        "严重度": severity,
                        "_sev_rank": sev_rank,
                        "_date": today["日期"],
                        "_deviation": deviation,
                    })

    if not alerts:
        return pd.DataFrame()
    out = pd.DataFrame(alerts).sort_values(
        ["_sev_rank", "_date", "_deviation"],
        ascending=[False, False, True],
    )
    return out


# ── 大盘对比（门店 vs 平台基准）─────────────────────────────────────────────
def compute_peer_comparison(df: pd.DataFrame) -> tuple[dict, pd.DataFrame]:
    """
    返回:
      - benchmarks: {平台: {下单转化率, 进店转化率, 客单价}}（加权大盘值）
      - detail: 每个门店在每个指标上与所在平台大盘的偏离
    """
    benchmarks = {}
    rows = []

    for plat in df["平台"].unique():
        sub = df[df["平台"] == plat]
        T_ord = sub["有效订单"].sum()
        T_vis = sub["进店人数"].sum()
        T_exp = sub["曝光人数"].sum()
        T_rev = sub["收入"].sum()

        b_order_conv = (T_ord / T_vis) if T_vis else 0
        b_visit_conv = (T_vis / T_exp) if T_exp else 0
        b_aov        = (T_rev / T_ord) if T_ord else 0

        benchmarks[plat] = {
            "下单转化率": b_order_conv,
            "进店转化率": b_visit_conv,
            "客单价":     b_aov,
            "_order_total": T_ord,
            "_visit_total": T_vis,
            "_exp_total":   T_exp,
            "_rev_total":   T_rev,
        }

        for store, grp in sub.groupby("门店名称"):
            s_ord = grp["有效订单"].sum()
            s_vis = grp["进店人数"].sum()
            s_exp = grp["曝光人数"].sum()
            s_rev = grp["收入"].sum()

            store_order_conv = (s_ord / s_vis) if s_vis else 0
            store_visit_conv = (s_vis / s_exp) if s_exp else 0
            store_aov        = (s_rev / s_ord) if s_ord else 0

            for name, sval, bval, is_pct in [
                ("下单转化率", store_order_conv, b_order_conv, True),
                ("进店转化率", store_visit_conv, b_visit_conv, True),
                ("客单价",     store_aov,        b_aov,        False),
            ]:
                if bval == 0:
                    continue
                dev = (sval - bval) / bval
                rows.append({
                    "门店": store, "平台": plat, "指标": name,
                    "门店值": sval, "大盘均值": bval, "偏离": dev,
                    "_is_pct": is_pct,
                })

    detail = pd.DataFrame(rows)
    return benchmarks, detail


def classify_peer_deviation(dev: float):
    """根据偏离度返回 (标签, 颜色) - 偏低坏，偏高好"""
    if dev <= -0.30:  return "🔴 严重偏低", "#dc2626"
    if dev <= -0.15:  return "🟠 偏低",     "#f59e0b"
    if dev >=  0.30:  return "🟢 标杆门店", "#10b981"
    if dev >=  0.15:  return "🔵 表现良好", "#3b82f6"
    return "⚪ 接近大盘", "#64748b"


# ── 利润分析 ──────────────────────────────────────────────────────────────────
def compute_profit_stats(df: pd.DataFrame) -> pd.DataFrame:
    """
    按 (门店, 平台) 聚合利润指标：
      - 实收(净收入)、客户实付、商家补贴投入
      - 单均实收、单均补贴
      - 实收率 = 实收 / 客户实付
      - 补贴 ROI = 实收 / 商家补贴投入
      - 补贴占比 = 补贴 / 客户实付
      - 优惠力度 = (优惠前总额 - 客户实付) / 优惠前总额
    """
    g = df.groupby(["门店名称", "平台"]).agg(
        实收=("收入", "sum"),
        客户实付=("顾客实付总额", "sum"),
        商家补贴=("商家补贴", "sum"),
        平台费=("平台费", "sum"),
        优惠前=("优惠前总额", "sum"),
        订单=("有效订单", "sum"),
    ).reset_index()

    g["单均实收"] = g["实收"] / g["订单"].replace(0, np.nan)
    g["单均补贴"] = g["商家补贴"] / g["订单"].replace(0, np.nan)
    g["实收率"]   = g["实收"] / g["客户实付"].replace(0, np.nan)
    g["补贴ROI"]  = g["实收"] / g["商家补贴"].replace(0, np.nan)
    g["补贴占比"] = g["商家补贴"] / g["客户实付"].replace(0, np.nan)
    g["优惠力度"] = (g["优惠前"] - g["客户实付"]) / g["优惠前"].replace(0, np.nan)

    return g.fillna(0)


# ── 流量类型退化检测（仅美团有此字段）────────────────────────────────────────
TRAFFIC_TYPE_TIER = {
    "高曝高转": 4, "高曝低转": 3,
    "低曝高转": 2, "低曝低转": 1,
}

def compute_traffic_type_alerts(df: pd.DataFrame) -> pd.DataFrame:
    """检测流量类型变化：当日 vs 7日、7日 vs 30日 退化"""
    sub = df[df["平台"] == "美团"].copy()
    if sub.empty or "流量类型_当日" not in sub.columns:
        return pd.DataFrame()

    sub = sub.dropna(subset=["流量类型_当日"])
    sub = sub[sub["流量类型_当日"].isin(TRAFFIC_TYPE_TIER.keys())]
    if sub.empty:
        return pd.DataFrame()

    rows = []
    for _, r in sub.iterrows():
        today  = r.get("流量类型_当日")
        d7     = r.get("流量类型_7日")
        d30    = r.get("流量类型_30日")
        t_today = TRAFFIC_TYPE_TIER.get(today, 0)
        t_7     = TRAFFIC_TYPE_TIER.get(d7, 0)
        t_30    = TRAFFIC_TYPE_TIER.get(d30, 0)

        events = []
        # 当日 vs 7日
        if t_today and t_7 and t_today < t_7:
            gap = t_7 - t_today
            sev = "🔴 严重" if gap >= 2 else "🟠 警告"
            events.append(("当日 vs 7日", f"{d7} → {today}", sev, gap))
        # 7日 vs 30日（更长期退化）
        if t_7 and t_30 and t_7 < t_30:
            gap = t_30 - t_7
            sev = "🔴 严重" if gap >= 2 else "🟠 警告"
            events.append(("7日 vs 30日", f"{d30} → {d7}", sev, gap))
        # 长期低迷
        if t_today == 1 and t_7 == 1 and t_30 == 1:
            events.append(("长期低迷", "低曝低转·持续", "🔴 严重", 0))

        for kind, change, sev, gap in events:
            rows.append({
                "日期": r["日期"].date(),
                "门店": r["门店名称"],
                "平台": "美团",
                "类型": kind,
                "变化": change,
                "当日": today, "7日": d7, "30日": d30,
                "严重度": sev,
                "_sev_rank": 3 if "严重" in sev else 2,
                "_date": r["日期"],
            })

    if not rows:
        return pd.DataFrame()
    return pd.DataFrame(rows).sort_values(["_sev_rank", "_date"], ascending=[False, False])


# ── 退单率异常检测（仅饿了么有此字段）────────────────────────────────────────
def compute_refund_alerts(df: pd.DataFrame, merchant_thresh=0.05, total_thresh=0.10) -> pd.DataFrame:
    """检测退单率异常：商户责任退单 > merchant_thresh 或总退单率 > total_thresh"""
    sub = df[df["平台"] == "饿了么"].copy()
    if sub.empty or "无效订单" not in sub.columns:
        return pd.DataFrame()

    sub["总单量"] = sub["有效订单"] + sub["无效订单"]
    sub = sub[sub["总单量"] > 0]
    if sub.empty:
        return pd.DataFrame()

    sub["总退单率"]     = sub["无效订单"] / sub["总单量"]
    sub["商户退单率"] = sub["商户原因无效订单"] / sub["总单量"]

    rows = []
    for _, r in sub.iterrows():
        m_rate = r["商户退单率"]
        t_rate = r["总退单率"]
        triggers = []
        if m_rate >= merchant_thresh:
            triggers.append(("商户责任退单率超标", m_rate, merchant_thresh, "商户原因"))
        if t_rate >= total_thresh and r["无效订单"] >= 3:
            triggers.append(("总退单率偏高", t_rate, total_thresh, "总退单"))

        for kind, val, thresh, sub_label in triggers:
            if val >= thresh * 2:    sev = "🔴 严重"
            elif val >= thresh * 1.5: sev = "🟠 警告"
            else:                     sev = "🟡 注意"
            rows.append({
                "日期": r["日期"].date(),
                "门店": r["门店名称"],
                "平台": "饿了么",
                "类型": kind,
                "退单率": f"{val*100:.2f}%",
                "阈值": f"{thresh*100:.0f}%",
                "无效订单数": int(r["无效订单"]),
                "商户原因数": int(r["商户原因无效订单"]),
                "严重度": sev,
                "_sev_rank": 3 if "严重" in sev else (2 if "警告" in sev else 1),
                "_date": r["日期"],
                "_value": val,
            })

    if not rows:
        return pd.DataFrame()
    return pd.DataFrame(rows).sort_values(["_sev_rank", "_value"], ascending=[False, False])


# ── 综合健康度评分 ────────────────────────────────────────────────────────────
def compute_health_score(df: pd.DataFrame) -> pd.DataFrame:
    """
    每家店 (门店, 平台) 综合健康度 0-100 分：
      营收规模 20 + 下单转化率 20 + 进店转化率 10 + 流量类型 10
      + 利润率 15 + 补贴ROI 10 + 退单率 10 + 新客占比合理度 5
    各项按平台内 min-max 或 vs 大盘标定，缺失项自动跳过并按权重补偿。
    """
    if df.empty:
        return pd.DataFrame()

    bench, _ = compute_peer_comparison(df)

    rows = []
    for (store, plat), grp in df.groupby(["门店名称", "平台"]):
        b = bench.get(plat, {})
        details = {}
        weights = {}

        # 营收规模 (vs 平台 top 1)
        plat_max_rev = df[df["平台"] == plat].groupby("门店名称")["收入"].sum().max()
        s_rev = grp["收入"].sum()
        if plat_max_rev and plat_max_rev > 0:
            details["营收规模"] = (s_rev / plat_max_rev) * 100
            weights["营收规模"] = 20

        # 下单转化率 (vs 大盘，截顶 150%)
        s_ord = grp["有效订单"].sum()
        s_vis = grp["进店人数"].sum()
        if s_vis > 0 and b.get("下单转化率", 0) > 0:
            ratio = (s_ord / s_vis) / b["下单转化率"]
            details["下单转化"] = min(ratio, 1.5) / 1.5 * 100
            weights["下单转化"] = 20

        # 进店转化率
        s_exp = grp["曝光人数"].sum()
        if s_exp > 0 and b.get("进店转化率", 0) > 0:
            ratio = (s_vis / s_exp) / b["进店转化率"]
            details["进店转化"] = min(ratio, 1.5) / 1.5 * 100
            weights["进店转化"] = 10

        # 流量类型（仅美团）
        if plat == "美团" and "流量类型_7日" in grp.columns:
            latest = grp.sort_values("日期").iloc[-1].get("流量类型_7日")
            if latest in TRAFFIC_TYPE_TIER:
                details["流量类型"] = (TRAFFIC_TYPE_TIER[latest] / 4) * 100
                weights["流量类型"] = 10

        # 利润率（实收/客付）vs 大盘
        s_paid = grp["顾客实付总额"].sum()
        if s_paid > 0:
            store_profit_rate = s_rev / s_paid
            plat_paid = df[df["平台"] == plat]["顾客实付总额"].sum()
            plat_rev  = df[df["平台"] == plat]["收入"].sum()
            bench_pr  = plat_rev / plat_paid if plat_paid else 0
            if bench_pr > 0:
                details["利润率"] = min(store_profit_rate / bench_pr, 1.5) / 1.5 * 100
                weights["利润率"] = 15

        # 补贴 ROI（实收/补贴）vs 大盘
        s_sub = grp["商家补贴"].sum()
        if s_sub > 0:
            store_roi = s_rev / s_sub
            plat_sub = df[df["平台"] == plat]["商家补贴"].sum()
            plat_rev2 = df[df["平台"] == plat]["收入"].sum()
            bench_roi = plat_rev2 / plat_sub if plat_sub else 0
            if bench_roi > 0:
                details["补贴ROI"] = min(store_roi / bench_roi, 1.5) / 1.5 * 100
                weights["补贴ROI"] = 10

        # 退单率（饿了么）— 越低越好
        if plat == "饿了么":
            s_inv = grp["无效订单"].sum()
            total = s_ord + s_inv
            if total > 0:
                refund_rate = s_inv / total
                # 0% = 100分, 10%+ = 0分
                details["退单率"] = max(0, (1 - refund_rate / 0.10)) * 100
                weights["退单率"] = 10

        # 新客占比合理度（理想 25%~45%）
        s_new = grp["新客进店人数"].sum() if "新客进店人数" in grp.columns else 0
        s_old = grp["老客进店人数"].sum() if "老客进店人数" in grp.columns else 0
        if (s_new + s_old) > 0:
            new_ratio = s_new / (s_new + s_old)
            if 0.25 <= new_ratio <= 0.45:
                details["新客占比"] = 100
            elif new_ratio < 0.25:
                details["新客占比"] = max(0, new_ratio / 0.25 * 100)
            else:
                details["新客占比"] = max(0, (1 - (new_ratio - 0.45) / 0.55) * 100)
            weights["新客占比"] = 5

        # 加权计算（缺项按存在项的权重归一）
        if not weights:
            continue
        total_w = sum(weights.values())
        score = sum(details[k] * weights[k] for k in weights) / total_w

        if   score >= 85: level, color = "🏆 标杆", "#10b981"
        elif score >= 70: level, color = "✅ 优秀", "#3b82f6"
        elif score >= 55: level, color = "⚪ 健康", "#64748b"
        elif score >= 40: level, color = "⚠️ 关注", "#f59e0b"
        else:             level, color = "🔴 危险", "#dc2626"

        rows.append({
            "门店": store, "平台": plat,
            "综合评分": round(score, 1),
            "等级": level,
            **{k: round(details[k], 1) for k in details},
            "_color": color,
        })

    if not rows:
        return pd.DataFrame()
    return pd.DataFrame(rows).sort_values("综合评分", ascending=False)


# ── 到手率分析：数据处理函数 ──────────────────────────────────────────────────────

def _rn(s):
    return str(s).strip().replace('　', '').replace('\xa0', '')

def _fc(df, *cands):
    m = {_rn(c): c for c in df.columns}
    for cand in cands:
        cn = _rn(cand)
        if cn in m:
            return m[cn]
        for k, v in m.items():
            if cn in k:
                return v
    return None

def _to_num(s):
    return pd.to_numeric(s, errors='coerce').fillna(0)

def rates_load_main(f):
    df = pd.read_excel(f, engine='openpyxl')
    df.columns = [_rn(c) for c in df.columns]
    c2, cn, cs, cp = _fc(df, '末级分类'), _fc(df, '商品名称'), _fc(df, '商品规格'), _fc(df, '门店售卖价格')
    miss = [n for n, c in [('末级分类', c2), ('商品名称', cn), ('门店售卖价格', cp)] if not c]
    if miss:
        raise ValueError(f"主表缺少必要列：{', '.join(miss)}")
    has_spec = cs is not None
    pdict = {}
    for _, row in df.iterrows():
        cat2, name = _rn(str(row[c2])), _rn(str(row[cn]))
        spec = _rn(str(row[cs])) if has_spec else ''
        v = row[cp]
        if pd.isna(v):
            continue
        try:
            price = float(v)
        except (ValueError, TypeError):
            continue
        pdict[(cat2, name, spec) if has_spec else (cat2, name)] = price
    return pdict, has_spec, len(pdict)

def _rates_lookup(pdict, has_spec, cat2, name, spec):
    c2n, nn, sn = _rn(str(cat2)), _rn(str(name)), _rn(str(spec))
    if has_spec:
        for k in [(c2n, nn, sn), (c2n, nn, '')]:
            if k in pdict:
                return pdict[k]
        for k, v in pdict.items():
            if k[0] == c2n and k[1] == nn:
                return v
    else:
        k = (c2n, nn)
        if k in pdict:
            return pdict[k]
    return None

def rates_load_channel(f, pdict, has_spec, ch_name):
    df = pd.read_excel(f, engine='openpyxl')
    df.columns = [_rn(c) for c in df.columns]
    c1 = _fc(df, '一级分类'); c2 = _fc(df, '末级分类'); cn = _fc(df, '商品名称')
    cs = _fc(df, '商品规格'); cq = _fc(df, '销售数量'); co = _fc(df, '销售原价')
    ca = _fc(df, '销售金额'); cd = _fc(df, '折扣金额')
    if c1:
        df = df[df[c1].astype(str).apply(_rn) == '成品'].copy()
    df = df.reset_index(drop=True)
    for col in [cq, co, ca, cd]:
        if col:
            df[col] = _to_num(df[col])
    rows, unmatched = [], []
    for _, row in df.iterrows():
        cat2 = str(row.get(c2, '')); name = str(row.get(cn, ''))
        spec = str(row.get(cs, '')) if cs else ''
        qty = float(row[cq]) if cq else 0
        orig = float(row[co]) if co else 0
        sale = float(row[ca]) if ca else 0
        disc = float(row[cd]) if cd else 0
        door_price = _rates_lookup(pdict, has_spec, cat2, name, spec)
        if door_price is not None:
            real = door_price * qty
            sys_r = sale / orig if orig > 0 else None
            real_r = sale / real if real > 0 else None
            prem = (orig - real) / real if real > 0 else None
            ok = True
        else:
            real = sys_r = real_r = prem = None; ok = False
            unmatched.append({'末级分类': _rn(cat2), '商品名称': _rn(name), '商品规格': _rn(spec)})
        rows.append({'末级分类': _rn(cat2), '商品名称': _rn(name), '商品规格': _rn(spec),
                     '销售数量': qty, '门店真实价': door_price, '销售原价': orig,
                     '门店真实销售额': real, '实收金额': sale, '折扣金额': disc,
                     '系统到手率': sys_r, '真实到手率': real_r, '溢价率': prem,
                     '匹配': '✓' if ok else '✗'})
    rdf = pd.DataFrame(rows)
    tot_o = rdf['销售原价'].sum(); tot_s = rdf['实收金额'].sum()
    tot_r = rdf['门店真实销售额'].dropna().sum(); tot_d = rdf['折扣金额'].sum()
    n = len(rdf); nm = n - len(unmatched)
    return rdf, {
        'sku_count': n, 'matched_count': nm, 'unmatched': unmatched,
        '销售数量': rdf['销售数量'].sum(), '系统销售原价': tot_o,
        '门店真实销售额': tot_r, '实收金额': tot_s, '折扣金额': tot_d,
        '系统到手率': tot_s / tot_o if tot_o > 0 else None,
        '真实到手率': tot_s / tot_r if tot_r > 0 else None,
        '溢价率': (tot_o - tot_r) / tot_r if tot_r > 0 else None,
    }

def rates_load_mp(f, delivery_fee):
    df = pd.read_excel(f, engine='openpyxl')
    df.columns = [_rn(c) for c in df.columns]
    c_id = _fc(df, '订单流水号'); c_mt = _fc(df, '用餐方式'); c_ch = _fc(df, '订单渠道')
    c_orig = _fc(df, '销售原价'); c_sale = _fc(df, '销售金额'); c_disc = _fc(df, '折扣金额')
    c_prof = _fc(df, '销售利润'); c_dn = _fc(df, '折扣名称'); c_pay = _fc(df, '支付方式')
    c_st = _fc(df, '门店名称'); c_tm = _fc(df, '订单时间')
    removed = 0
    if c_id:
        mask = df[c_id].astype(str).apply(_rn) == '汇总'
        removed = int(mask.sum()); df = df[~mask].copy()
    if c_mt:
        df = df[df[c_mt].astype(str).apply(_rn) == '外卖'].copy()
    df = df.reset_index(drop=True)
    warn = None
    if c_ch and len(df) > 0:
        bad = (df[c_ch].astype(str).apply(_rn) != '自营小程序').sum()
        if bad > 0:
            warn = f"⚠️ {bad} 行订单渠道不是「自营小程序」"
    for col in [c_orig, c_sale, c_disc, c_prof]:
        if col:
            df[col] = _to_num(df[col])
    ot = df[c_orig].sum() if c_orig else 0
    st_ = df[c_sale].sum() if c_sale else 0
    dt = df[c_disc].sum() if c_disc else 0
    pt = df[c_prof].sum() if c_prof else 0
    net = st_ - delivery_fee
    cols = {}
    if c_id: cols['订单流水号'] = df[c_id]
    if c_st: cols['门店名称'] = df[c_st]
    if c_mt: cols['用餐方式'] = df[c_mt]
    if c_tm: cols['订单时间'] = df[c_tm]
    if c_orig: cols['销售原价'] = df[c_orig]
    if c_sale: cols['实收金额'] = df[c_sale]
    if c_disc: cols['折扣金额'] = df[c_disc]
    if c_dn: cols['折扣名称'] = df[c_dn]
    if c_pay: cols['支付方式'] = df[c_pay]
    return pd.DataFrame(cols) if cols else pd.DataFrame(), {
        'order_count': len(df), 'removed': removed, 'warning': warn,
        '系统销售原价': ot, '门店真实销售额': ot,
        '实收金额(毛)': st_, '配送费': delivery_fee, '净实收': net,
        '折扣金额': dt, '利润合计': pt,
        '毛到手率': st_ / ot if ot > 0 else None,
        '净到手率': net / ot if ot > 0 else None,
        '实收金额': net,
        '系统到手率': net / ot if ot > 0 else None,
        '真实到手率': net / ot if ot > 0 else None,
        '溢价率': 0.0,
    }


def generate_daily_report(df, report_date) -> pd.DataFrame:
    day = df[df["日期"] == pd.Timestamp(report_date)]
    if day.empty:
        return pd.DataFrame()
    summary = day.groupby(["门店名称", "平台"]).agg(
        营业收入=("收入", "sum"),
        有效订单=("有效订单", "sum"),
        单均实付=("单均实付", "mean"),
        曝光人数=("曝光人数", "sum"),
        进店人数=("进店人数", "sum"),
        进店转化率=("进店转化率", "mean"),
        下单转化率=("下单转化率", "mean"),
    ).reset_index().sort_values("营业收入", ascending=False)
    summary["营业收入"] = summary["营业收入"].map(lambda x: f"¥{x:,.2f}")
    summary["单均实付"] = summary["单均实付"].map(lambda x: f"¥{x:.2f}")
    summary["进店转化率"] = summary["进店转化率"].map(lambda x: f"{x*100:.1f}%")
    summary["下单转化率"] = summary["下单转化率"].map(lambda x: f"{x*100:.1f}%")
    summary["曝光人数"] = summary["曝光人数"].map(lambda x: f"{int(x):,}")
    summary["进店人数"] = summary["进店人数"].map(lambda x: f"{int(x):,}")
    summary["有效订单"] = summary["有效订单"].map(lambda x: f"{int(x):,}")
    return summary.rename(columns={"门店名称": "门店"})


# ── PPT 日报生成 ──────────────────────────────────────────────────────────────
PPT_FONT = "微软雅黑"
PPT_COLOR_PRIMARY  = RGBColor(0x4F, 0x46, 0xE5)  # 主色 紫
PPT_COLOR_SUCCESS  = RGBColor(0x10, 0xB9, 0x81)  # 绿
PPT_COLOR_WARNING  = RGBColor(0xF5, 0x9E, 0x0B)  # 橙
PPT_COLOR_DANGER   = RGBColor(0xDC, 0x26, 0x26)  # 红
PPT_COLOR_TEXT     = RGBColor(0x0F, 0x17, 0x2A)  # 主文字
PPT_COLOR_SUBTEXT  = RGBColor(0x64, 0x74, 0x8B)  # 次要文字
PPT_COLOR_LIGHTBG  = RGBColor(0xF8, 0xFA, 0xFC)  # 浅底
PPT_COLOR_BORDER   = RGBColor(0xE2, 0xE8, 0xF0)  # 边框


def _ppt_set_text(tf, text, size=14, bold=False, color=PPT_COLOR_TEXT,
                  align=PP_ALIGN.LEFT, font=PPT_FONT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _ppt_textbox(slide, x, y, w, h, text, **kwargs):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    _ppt_set_text(tf, text, **kwargs)
    return box


def _ppt_kpi_card(slide, x, y, w, h, label, value, accent=PPT_COLOR_PRIMARY):
    """画一个 KPI 卡片：左边竖条 + 标签 + 大数字"""
    # 卡片背景
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.color.rgb = PPT_COLOR_BORDER
    bg.line.width = Pt(0.75)
    bg.shadow.inherit = False

    # 左侧色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # 标签
    _ppt_textbox(slide, x + 0.18, y + 0.12, w - 0.25, 0.3,
                 label, size=11, color=PPT_COLOR_SUBTEXT)
    # 数值
    _ppt_textbox(slide, x + 0.18, y + 0.42, w - 0.25, h - 0.5,
                 value, size=22, bold=True, color=PPT_COLOR_TEXT)


def _ppt_table(slide, x, y, w, h, headers, data, header_color=PPT_COLOR_PRIMARY):
    """渲染一个简洁表格"""
    rows = len(data) + 1
    cols = len(headers)
    if rows < 2:
        return
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(x), Inches(y),
                                          Inches(w), Inches(h))
    table = table_shape.table

    # 表头
    for ci, h_text in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        _ppt_set_text(tf, h_text, size=11, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # 数据行
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (PPT_COLOR_LIGHTBG if ri % 2 == 0
                                         else RGBColor(0xFF, 0xFF, 0xFF))
            tf = cell.text_frame
            align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            _ppt_set_text(tf, val, size=10, color=PPT_COLOR_TEXT, align=align)


def _ppt_add_slide_header(slide, title, subtitle=""):
    """每页顶部统一的页头"""
    # 顶部色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPT_COLOR_PRIMARY
    bar.line.fill.background()

    _ppt_textbox(slide, 0.4, 0.2, 12, 0.5,
                 title, size=22, bold=True, color=PPT_COLOR_TEXT)
    if subtitle:
        _ppt_textbox(slide, 0.4, 0.7, 12, 0.4,
                     subtitle, size=12, color=PPT_COLOR_SUBTEXT)


def generate_ppt_report(df_all, report_date,
                         alerts_df=None, traffic_alerts=None, refund_alerts=None) -> io.BytesIO:
    """生成日报 PPT，返回 BytesIO。"""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    day = df_all[df_all["日期"] == pd.Timestamp(report_date)]
    if day.empty:
        # 空数据兜底
        slide = prs.slides.add_slide(blank)
        _ppt_textbox(slide, 1, 3, 11, 1, f"{report_date} 无数据",
                     size=28, bold=True, color=PPT_COLOR_DANGER, align=PP_ALIGN.CENTER)
        buf = io.BytesIO()
        prs.save(buf); buf.seek(0); return buf

    # ── 数据汇总 ───────────────────────────────────
    total_rev = day["收入"].sum()
    total_ord = day["有效订单"].sum()
    total_exp = day["曝光人数"].sum() if "曝光人数" in day.columns else 0
    total_vis = day["进店人数"].sum() if "进店人数" in day.columns else 0
    avg_price = (total_rev / total_ord) if total_ord else 0
    overall_conv = (total_ord / total_vis) if total_vis else 0
    store_cnt = day["门店名称"].nunique()

    plat_summary = day.groupby("平台").agg(
        营收=("收入", "sum"),
        订单=("有效订单", "sum"),
        曝光=("曝光人数", "sum"),
        进店=("进店人数", "sum"),
    ).reset_index()
    plat_summary["客单价"] = plat_summary["营收"] / plat_summary["订单"].replace(0, np.nan)
    plat_summary["下单转化"] = plat_summary["订单"] / plat_summary["进店"].replace(0, np.nan)

    store_rank = day.groupby(["门店名称", "平台"]).agg(
        营收=("收入", "sum"),
        订单=("有效订单", "sum"),
    ).reset_index().sort_values("营收", ascending=False)

    # ── Slide 1：封面 ──────────────────────────────
    slide = prs.slides.add_slide(blank)
    # 整页彩色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPT_COLOR_PRIMARY
    bg.line.fill.background()

    _ppt_textbox(slide, 0.8, 1.8, 12, 1,
                 "门店运营日报", size=44, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    _ppt_textbox(slide, 0.8, 2.8, 12, 0.6,
                 f"📅 {report_date}", size=20,
                 color=RGBColor(0xE0, 0xE7, 0xFF))

    # 三大核心数字
    _ppt_textbox(slide, 0.8, 4.0, 12, 0.4,
                 "今日核心数据", size=14,
                 color=RGBColor(0xC7, 0xD2, 0xFE))

    def _cover_metric(x, label, value):
        _ppt_textbox(slide, x, 4.5, 4, 0.4, label, size=13,
                      color=RGBColor(0xC7, 0xD2, 0xFE))
        _ppt_textbox(slide, x, 4.9, 4, 0.9, value, size=36, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))

    _cover_metric(0.8, "总营业收入", fmt_money(total_rev))
    _cover_metric(5.0, "总有效订单", f"{int(total_ord):,} 单")
    _cover_metric(9.2, "活跃门店",   f"{store_cnt} 家")

    _ppt_textbox(slide, 0.8, 6.8, 12, 0.4,
                 "本报告由门店运营监控看板自动生成", size=10,
                 color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.LEFT)

    # ── Slide 2：核心指标 + 平台对比 ────────────────
    slide = prs.slides.add_slide(blank)
    _ppt_add_slide_header(slide, "📊 核心指标概览",
                           f"{report_date} | {store_cnt} 家活跃门店")

    # 5 个 KPI 卡
    card_w = 2.4; card_h = 1.4; gap = 0.15
    start_x = 0.4; start_y = 1.4
    for i, (label, value, accent) in enumerate([
        ("总营业收入", fmt_money(total_rev), PPT_COLOR_PRIMARY),
        ("总有效订单", f"{int(total_ord):,}", PPT_COLOR_SUCCESS),
        ("加权客单价", f"¥{avg_price:.2f}", PPT_COLOR_PRIMARY),
        ("总曝光人数", f"{int(total_exp):,}", PPT_COLOR_WARNING),
        ("整体下单转化", f"{overall_conv*100:.1f}%", PPT_COLOR_SUCCESS),
    ]):
        x = start_x + i * (card_w + gap)
        _ppt_kpi_card(slide, x, start_y, card_w, card_h, label, value, accent)

    # 平台对比表
    _ppt_textbox(slide, 0.4, 3.2, 12, 0.4,
                 "平台拆分", size=14, bold=True, color=PPT_COLOR_TEXT)

    headers = ["平台", "营收", "订单", "客单价", "曝光人数", "进店人数", "下单转化"]
    data = []
    for _, r in plat_summary.iterrows():
        data.append([
            r["平台"],
            f"¥{r['营收']:,.2f}",
            f"{int(r['订单']):,}",
            f"¥{r['客单价']:.2f}" if pd.notna(r['客单价']) else "—",
            f"{int(r['曝光']):,}",
            f"{int(r['进店']):,}",
            f"{r['下单转化']*100:.1f}%" if pd.notna(r['下单转化']) else "—",
        ])
    _ppt_table(slide, 0.4, 3.7, 12.5, 1.6, headers, data)

    # ── Slide 3：Top 5 门店 ─────────────────────────
    slide = prs.slides.add_slide(blank)
    _ppt_add_slide_header(slide, "🏆 营收 Top 5 门店",
                           "今日表现最好的门店")

    top5 = store_rank.head(5)
    headers = ["排名", "门店", "平台", "营收", "订单"]
    data = []
    for i, (_, r) in enumerate(top5.iterrows(), 1):
        data.append([
            f"#{i}",
            r["门店名称"],
            r["平台"],
            f"¥{r['营收']:,.2f}",
            f"{int(r['订单']):,} 单",
        ])
    if data:
        _ppt_table(slide, 0.4, 1.6, 12.5, 4, headers, data,
                    header_color=PPT_COLOR_SUCCESS)
    _ppt_textbox(slide, 0.4, 6.5, 12, 0.5,
                 f"📌 Top 5 合计营收 ¥{top5['营收'].sum():,.2f}（占总营收 {top5['营收'].sum()/total_rev*100 if total_rev else 0:.1f}%）",
                 size=12, color=PPT_COLOR_SUBTEXT)

    # ── Slide 4：需关注门店 ─────────────────────────
    slide = prs.slides.add_slide(blank)
    _ppt_add_slide_header(slide, "⚠️ 需关注门店",
                           "今日营收末位 + 异常门店")

    bottom5 = store_rank.tail(5).sort_values("营收")
    headers = ["排名", "门店", "平台", "营收", "订单"]
    data = []
    total_rank = len(store_rank)
    for i, (_, r) in enumerate(bottom5.iterrows()):
        rank_idx = total_rank - len(bottom5) + i + 1
        data.append([
            f"#{rank_idx}",
            r["门店名称"],
            r["平台"],
            f"¥{r['营收']:,.2f}",
            f"{int(r['订单']):,} 单",
        ])
    if data:
        _ppt_table(slide, 0.4, 1.6, 12.5, 3.5, headers, data,
                    header_color=PPT_COLOR_WARNING)

    # ── Slide 5：预警汇总 ──────────────────────────
    slide = prs.slides.add_slide(blank)
    _ppt_add_slide_header(slide, "🚨 当日预警汇总",
                           "指标异常 / 流量类型退化 / 退单率异常")

    # 当日相关预警过滤
    def _filter_today(df_, date_col="_date"):
        if df_ is None or df_.empty or date_col not in df_.columns:
            return df_ if df_ is not None else pd.DataFrame()
        return df_[df_[date_col].dt.date == report_date]

    a_today = _filter_today(alerts_df)
    t_today = _filter_today(traffic_alerts)
    r_today = _filter_today(refund_alerts)

    # 3 个汇总卡片
    sev_a = (a_today["严重度"].str.contains("严重").sum() if not a_today.empty else 0)
    sev_t = (t_today["严重度"].str.contains("严重").sum() if not t_today.empty else 0)
    sev_r = (r_today["严重度"].str.contains("严重").sum() if not r_today.empty else 0)

    _ppt_kpi_card(slide, 0.4, 1.4, 4.1, 1.2, "指标异常",
                   f"{len(a_today)} 条 / 严重 {sev_a}", PPT_COLOR_DANGER)
    _ppt_kpi_card(slide, 4.7, 1.4, 4.1, 1.2, "流量类型退化",
                   f"{len(t_today)} 条 / 严重 {sev_t}", PPT_COLOR_WARNING)
    _ppt_kpi_card(slide, 9.0, 1.4, 4.1, 1.2, "退单率异常",
                   f"{len(r_today)} 条 / 严重 {sev_r}", PPT_COLOR_DANGER)

    # 严重预警明细表（合并三类，最多取 8 条）
    severe_rows = []
    for d, kind in [(a_today, "指标"), (t_today, "流量"), (r_today, "退单")]:
        if d is None or d.empty:
            continue
        severe = d[d["严重度"].str.contains("严重")] if "严重度" in d.columns else d
        for _, r in severe.iterrows():
            store = r.get("门店", "")
            plat  = r.get("平台", "")
            if kind == "指标":
                detail = f"{r.get('指标','')}：{r.get('当日值','')}（vs 历史 {r.get('历史均值','')}, {r.get('偏离','')}）"
            elif kind == "流量":
                detail = f"{r.get('类型','')}：{r.get('变化','')}"
            else:
                detail = f"{r.get('类型','')}：{r.get('退单率','')}（{int(r.get('无效订单数',0))} 单无效）"
            severe_rows.append([kind, store, plat, detail])

    if severe_rows:
        _ppt_textbox(slide, 0.4, 2.9, 12, 0.4,
                      "🔴 严重预警明细（最多展示 8 条）", size=14, bold=True,
                      color=PPT_COLOR_DANGER)
        _ppt_table(slide, 0.4, 3.4, 12.5, 3.5,
                    ["类别", "门店", "平台", "详情"],
                    severe_rows[:8],
                    header_color=PPT_COLOR_DANGER)
    else:
        _ppt_textbox(slide, 0.4, 3.5, 12, 1,
                      "✅ 今日未发现严重级别预警", size=20, bold=True,
                      color=PPT_COLOR_SUCCESS, align=PP_ALIGN.CENTER)

    # ── 保存 ───────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── 侧边栏 ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("## 🏪 门店运营看板")
    st.caption("美团 + 饿了么 数据监控")
    st.divider()

    st.markdown("#### 📁 数据文件")
    mt_upload  = st.file_uploader("美团 CSV", type=["csv"], key="mt")
    elm_upload = st.file_uploader("饿了么 Excel", type=["xlsx", "xls"], key="elm")

    # 本地调试时可选：从环境变量读取默认文件路径
    DEFAULT_MT  = os.environ.get("DASHBOARD_MT_FILE")
    DEFAULT_ELM = os.environ.get("DASHBOARD_ELM_FILE")

    mt_src  = mt_upload  if mt_upload  else (DEFAULT_MT  if DEFAULT_MT  and os.path.exists(DEFAULT_MT)  else None)
    elm_src = elm_upload if elm_upload else (DEFAULT_ELM if DEFAULT_ELM and os.path.exists(DEFAULT_ELM) else None)

    with st.spinner("加载数据..."):
        df_all = load_all(mt_src, elm_src)

    if df_all.empty:
        st.info("👆 请上传美团 CSV 和/或饿了么 Excel 文件开始使用")
        st.markdown("""
        <div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;
        padding:14px 18px;font-size:13px;color:#475569;margin-top:10px;'>
        <b style='color:#0f172a;'>📥 使用步骤：</b>
        <ol style='margin:6px 0 0 0;padding-left:20px;'>
        <li>从美团/饿了么商家后台导出门店数据</li>
        <li>点击上方上传按钮选择文件</li>
        <li>右侧自动生成看板和预警</li>
        </ol>
        </div>
        """, unsafe_allow_html=True)
        st.stop()

    st.success(f"✅ 已加载 {len(df_all):,} 条记录")

    st.divider()
    st.markdown("#### 🔍 筛选")

    all_dates = sorted(df_all["日期"].dt.date.unique())
    min_d, max_d = all_dates[0], all_dates[-1]

    # 健壮的日期范围处理（单选/双选都OK）
    raw_range = st.date_input(
        "日期范围",
        value=(min_d, max_d),
        min_value=min_d,
        max_value=max_d,
        key="date_range_input",
    )
    if isinstance(raw_range, (tuple, list)):
        if len(raw_range) == 2:
            start_d, end_d = raw_range[0], raw_range[1]
        elif len(raw_range) == 1:
            start_d = end_d = raw_range[0]
        else:
            start_d, end_d = min_d, max_d
    else:
        start_d = end_d = raw_range

    platforms = st.multiselect(
        "平台",
        df_all["平台"].unique().tolist(),
        default=df_all["平台"].unique().tolist(),
    )

    all_stores = sorted(df_all["门店名称"].unique().tolist())
    selected_stores = st.multiselect("门店（留空=全部）", all_stores)

    st.divider()
    st.markdown("#### ⚠️ 预警阈值")
    z_thresh = st.slider(
        "Z-score 阈值",
        min_value=1.0, max_value=3.0, value=2.0, step=0.1,
        help="今日值偏离历史均值多少个标准差时触发。建议 1.5~2.5。",
    )
    drop_thresh = st.slider(
        "环比下跌阈值",
        min_value=0.1, max_value=0.6, value=0.30, step=0.05,
        format="%.0f%%",
        help="对比历史均值，下跌超过该比例触发预警。",
    )
    monitor_metrics = st.multiselect(
        "监控指标",
        list(ALERT_METRICS.keys()),
        default=["营业收入", "进店转化率", "下单转化率", "曝光人数"],
    )


# ── 过滤数据 ──────────────────────────────────────────────────────────────────
if not platforms:
    st.warning("⚠️ 请至少选择一个平台")
    st.stop()

mask = (
    (df_all["日期"].dt.date >= start_d) &
    (df_all["日期"].dt.date <= end_d) &
    (df_all["平台"].isin(platforms))
)
if selected_stores:
    mask &= df_all["门店名称"].isin(selected_stores)
df = df_all[mask].copy()

if df.empty:
    st.warning("⚠️ 当前筛选条件下无数据，请调整筛选条件")
    st.stop()


# ── 顶部信息条 ────────────────────────────────────────────────────────────────
st.markdown(f"""
<h1 style='font-size:24px;margin-bottom:8px;color:#0f172a;'>📊 门店运营监控看板</h1>
<div class="top-info" style='font-size:13px;color:#475569;'>
    📅 周期 <b>{start_d} ~ {end_d}</b> &nbsp;&nbsp;|&nbsp;&nbsp;
    🏷 平台 <b>{"、".join(platforms)}</b> &nbsp;&nbsp;|&nbsp;&nbsp;
    🏪 门店数 <b>{df["门店名称"].nunique()} 家</b> &nbsp;&nbsp;|&nbsp;&nbsp;
    📦 记录数 <b>{len(df):,}</b>
</div>
""", unsafe_allow_html=True)

tab1, tab2, tab3, tab_profit, tab_peer, tab_health, tab_alert, tab4, tab_rates = st.tabs(
    ["📈 综合概览", "🏪 门店分析", "👥 流量分析",
     "💰 利润分析", "🎯 大盘对比", "⭐ 健康度",
     "⚠️ 预警监控", "📋 日报", "💸 到手率分析"]
)

# ── Tab 1: 综合概览 ───────────────────────────────────────────────────────────
with tab1:
    total_rev  = df["收入"].sum()
    total_ord  = df["有效订单"].sum()
    avg_price  = (total_rev / total_ord) if total_ord else 0
    total_expo = df["曝光人数"].sum() if "曝光人数" in df.columns else 0
    total_vis  = df["进店人数"].sum() if "进店人数" in df.columns else 0
    overall_conv = (total_ord / total_vis) if total_vis else 0

    c1, c2, c3, c4, c5 = st.columns(5)
    with c1: kpi_card("总营业收入", fmt_money(total_rev), f"原值 ¥{total_rev:,.2f}")
    with c2: kpi_card("总有效订单", fmt_int(total_ord), "单")
    with c3: kpi_card("加权客单价", f"¥{avg_price:.2f}", "= 总收入 / 总订单")
    with c4: kpi_card("总曝光人数", fmt_int(total_expo), "人")
    with c5: kpi_card("整体下单转化", fmt_pct(overall_conv), "= 订单 / 进店")

    st.markdown('<div class="section-title">每日营业收入 & 订单量趋势</div>', unsafe_allow_html=True)
    st.plotly_chart(plot_daily_trend(df), width='stretch')

    st.markdown('<div class="section-title">各平台汇总</div>', unsafe_allow_html=True)
    plat_summary = df.groupby("平台").agg(
        营业收入=("收入", "sum"),
        有效订单=("有效订单", "sum"),
        曝光人数=("曝光人数", "sum"),
        进店人数=("进店人数", "sum"),
    ).reset_index()
    plat_summary["客单价"] = plat_summary["营业收入"] / plat_summary["有效订单"].replace(0, 1)
    plat_summary["下单转化率"] = plat_summary["有效订单"] / plat_summary["进店人数"].replace(0, 1)
    plat_summary["营业收入"]  = plat_summary["营业收入"].map(lambda x: f"¥{x:,.2f}")
    plat_summary["有效订单"]  = plat_summary["有效订单"].map(lambda x: f"{int(x):,}")
    plat_summary["曝光人数"]  = plat_summary["曝光人数"].map(lambda x: f"{int(x):,}")
    plat_summary["进店人数"]  = plat_summary["进店人数"].map(lambda x: f"{int(x):,}")
    plat_summary["客单价"]    = plat_summary["客单价"].map(lambda x: f"¥{x:.2f}")
    plat_summary["下单转化率"]= plat_summary["下单转化率"].map(lambda x: f"{x*100:.1f}%")
    st.dataframe(plat_summary, width='stretch', hide_index=True)


# ── Tab 2: 门店分析 ───────────────────────────────────────────────────────────
with tab2:
    col_left, col_right = st.columns(2)
    with col_left:
        st.markdown('<div class="section-title">门店营业收入排名</div>', unsafe_allow_html=True)
        st.plotly_chart(plot_store_rank(df, "收入", "营业收入 (¥)"), width='stretch')
    with col_right:
        st.markdown('<div class="section-title">门店订单量排名</div>', unsafe_allow_html=True)
        st.plotly_chart(plot_store_rank(df, "有效订单", "有效订单 (单)"), width='stretch')

    st.markdown('<div class="section-title">门店明细数据</div>', unsafe_allow_html=True)
    store_detail = df.groupby(["门店名称", "平台"]).agg(
        营业收入=("收入", "sum"),
        有效订单=("有效订单", "sum"),
        曝光人数=("曝光人数", "sum"),
        进店人数=("进店人数", "sum"),
    ).reset_index().sort_values("营业收入", ascending=False)
    store_detail["客单价"] = store_detail["营业收入"] / store_detail["有效订单"].replace(0, 1)
    store_detail["下单转化率"] = store_detail["有效订单"] / store_detail["进店人数"].replace(0, 1)
    store_detail["营业收入"]   = store_detail["营业收入"].map(lambda x: f"¥{x:,.2f}")
    store_detail["有效订单"]   = store_detail["有效订单"].map(lambda x: f"{int(x):,}")
    store_detail["曝光人数"]   = store_detail["曝光人数"].map(lambda x: f"{int(x):,}")
    store_detail["进店人数"]   = store_detail["进店人数"].map(lambda x: f"{int(x):,}")
    store_detail["客单价"]     = store_detail["客单价"].map(lambda x: f"¥{x:.2f}")
    store_detail["下单转化率"] = store_detail["下单转化率"].map(lambda x: f"{x*100:.1f}%")
    st.dataframe(store_detail, width='stretch', hide_index=True)


# ── Tab 3: 流量分析（按平台拆分）──────────────────────────────────────────────
with tab3:
    platforms_in_df = df["平台"].unique().tolist()

    if len(platforms_in_df) == 0:
        st.info("当前筛选下无数据")
    else:
        # 每个平台一列
        cols = st.columns(len(platforms_in_df))
        for i, plat in enumerate(platforms_in_df):
            sub = df[df["平台"] == plat]
            with cols[i]:
                color = PLATFORM_COLORS.get(plat, "#888")
                st.markdown(f"""
<div style='background:{color}1A;border-left:4px solid {color};
padding:10px 16px;border-radius:6px;margin-bottom:8px;'>
<b style='color:#0f172a;font-size:16px;'>🏷 {plat}</b>
<span style='color:#64748b;font-size:12px;margin-left:10px;'>
{sub['门店名称'].nunique()} 家门店 · {len(sub):,} 条记录
</span>
</div>
""", unsafe_allow_html=True)

                # 平台级 KPI
                p_exp = sub["曝光人数"].sum()
                p_vis = sub["进店人数"].sum()
                p_ord = sub["有效订单"].sum()
                p_rev = sub["收入"].sum()
                k1, k2, k3 = st.columns(3)
                with k1: kpi_card("总营业额", fmt_money(p_rev))
                with k2: kpi_card("曝光→进店", f"{(p_vis/p_exp*100) if p_exp else 0:.2f}%")
                with k3: kpi_card("进店→下单", f"{(p_ord/p_vis*100) if p_vis else 0:.2f}%")

                st.markdown(f'<div class="section-title">{plat} · 漏斗</div>', unsafe_allow_html=True)
                st.plotly_chart(plot_funnel(sub), width='stretch', key=f"funnel_{plat}")

                st.markdown(f'<div class="section-title">{plat} · 新客 vs 老客</div>', unsafe_allow_html=True)
                st.plotly_chart(plot_new_old(sub), width='stretch', key=f"newold_{plat}")

                st.markdown(f'<div class="section-title">{plat} · 各门店转化率</div>', unsafe_allow_html=True)
                st.plotly_chart(plot_conversion_by_store(sub), width='stretch', key=f"conv_{plat}")


# ── Tab 健康度评分 ────────────────────────────────────────────────────────────
with tab_health:
    health_df = compute_health_score(df)
    if health_df.empty:
        st.info("当前筛选下无数据")
    else:
        # 顶部分布统计
        lvl_count = health_df["等级"].value_counts().to_dict()
        c1, c2, c3, c4, c5 = st.columns(5)
        for col, key, label in [
            (c1, "🏆 标杆", "🏆 标杆 (≥85)"),
            (c2, "✅ 优秀", "✅ 优秀 (70-85)"),
            (c3, "⚪ 健康", "⚪ 健康 (55-70)"),
            (c4, "⚠️ 关注", "⚠️ 关注 (40-55)"),
            (c5, "🔴 危险", "🔴 危险 (<40)"),
        ]:
            with col:
                kpi_card(label, str(lvl_count.get(key, 0)), "家")

        st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:12px 16px;
font-size:12px;color:#475569;margin:8px 0;'>
<b style='color:#0f172a;'>评分构成：</b>
营收规模 20 + 下单转化 20 + 进店转化 10 + 流量类型 10 +
利润率 15 + 补贴ROI 10 + 退单率 10 + 新客占比 5 = 100 分。
<br>各项均按平台内对标计算，缺项自动按权重补偿。
</div>
""", unsafe_allow_html=True)

        # 评分条形图
        st.markdown('<div class="section-title">门店健康度排行</div>', unsafe_allow_html=True)
        hdf = health_df.copy()
        hdf["门店短名"] = hdf["门店"].apply(_short_store_name) + " [" + hdf["平台"] + "]"
        hdf = hdf.sort_values("综合评分", ascending=True)
        fig = go.Figure(go.Bar(
            y=hdf["门店短名"],
            x=hdf["综合评分"],
            orientation="h",
            marker_color=hdf["_color"],
            text=[f"{v:.1f} {l}" for v, l in zip(hdf["综合评分"], hdf["等级"])],
            textposition="outside",
            textfont=dict(color=TEXT_COLOR, size=11),
        ))
        _base_layout(fig, height=max(400, len(hdf) * 28))
        fig.update_layout(margin=dict(l=10, r=120, t=30, b=10), showlegend=False)
        fig.update_xaxes(gridcolor=GRID_COLOR, range=[0, 110], showline=True, linecolor=GRID_COLOR)
        fig.update_yaxes(gridcolor=GRID_COLOR, automargin=True)
        # 加分级参考线
        for x_val, label, color in [(40, "危险线", "#dc2626"), (55, "关注线", "#f59e0b"),
                                      (70, "健康线", "#3b82f6"), (85, "标杆线", "#10b981")]:
            fig.add_vline(x=x_val, line=dict(color=color, width=1, dash="dot"),
                          annotation_text=label, annotation_position="top",
                          annotation_font_color=color, annotation_font_size=10)
        st.plotly_chart(fig, width='stretch')

        # 单店雷达图
        st.markdown('<div class="section-title">单店维度详情（雷达图）</div>', unsafe_allow_html=True)
        store_options = hdf.sort_values("综合评分", ascending=False)
        store_options["选项"] = store_options["综合评分"].map(lambda x: f"{x:.1f}分") + " · " + \
                              store_options["门店"] + " [" + store_options["平台"] + "]"
        opt_list = store_options["选项"].tolist()
        sel_idx = st.selectbox("选择门店", range(len(opt_list)),
                                format_func=lambda i: opt_list[i], key="health_sel")
        row = store_options.iloc[sel_idx]
        radar_cols = ["营收规模", "下单转化", "进店转化", "流量类型",
                      "利润率", "补贴ROI", "退单率", "新客占比"]
        # 只保留该店真正有得分的维度，labels 和 data 长度对齐
        radar_pairs = []
        for c in radar_cols:
            v = row.get(c)
            if v is not None and not pd.isna(v):
                radar_pairs.append((c, float(v)))
        if len(radar_pairs) < 3:
            st.info("该店可对比维度不足，无法绘制雷达图")
            radar_labels, radar_data = [], []
        else:
            radar_labels = [p[0] for p in radar_pairs]
            radar_data   = [p[1] for p in radar_pairs]

        col_l, col_r = st.columns([2, 1])
        store_color = str(row.get("_color") or "#64748b")
        if not store_color.startswith("#"):
            store_color = "#64748b"
        with col_l:
            if radar_data:
                try:
                    fig_radar = go.Figure(go.Scatterpolar(
                        r=radar_data + [radar_data[0]],
                        theta=radar_labels + [radar_labels[0]],
                        fill="toself",
                        fillcolor=hex_to_rgba(store_color, 0.25),
                        line=dict(color=store_color, width=2),
                        marker=dict(color=store_color, size=8),
                    ))
                    fig_radar.update_layout(
                        polar=dict(
                            radialaxis=dict(visible=True, range=[0, 100], gridcolor=GRID_COLOR,
                                            tickfont=dict(color=TEXT_COLOR, size=10)),
                            angularaxis=dict(tickfont=dict(color=TEXT_COLOR, size=12)),
                            bgcolor="#fff",
                        ),
                        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
                        font=dict(color=TEXT_COLOR),
                        height=400, margin=dict(l=40, r=40, t=40, b=40),
                        showlegend=False,
                    )
                    st.plotly_chart(fig_radar, width='stretch')
                except Exception as e:
                    st.warning(f"雷达图渲染失败：{e}")
        with col_r:
            st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-left:4px solid {store_color};
border-radius:10px;padding:16px 20px;margin-top:30px;'>
<div style='font-size:13px;color:#64748b;margin-bottom:4px;'>{row['平台']}</div>
<div style='font-size:16px;font-weight:700;color:#0f172a;margin-bottom:8px;'>{row['门店']}</div>
<div style='font-size:32px;font-weight:700;color:{store_color};margin-bottom:4px;'>{row['综合评分']:.1f}</div>
<div style='font-size:14px;color:{store_color};margin-bottom:12px;'>{row['等级']}</div>
<hr style='border:none;border-top:1px solid #e2e8f0;'>
<div style='font-size:12px;color:#64748b;'>各维度得分：</div>
""" + "".join([
                f"<div style='display:flex;justify-content:space-between;padding:3px 0;font-size:13px;'>"
                f"<span style='color:#475569;'>{c}</span>"
                f"<b style='color:#0f172a;'>{row[c]:.1f}</b></div>"
                for c in radar_labels
            ]) + "</div>", unsafe_allow_html=True)

        # 明细表
        st.markdown('<div class="section-title">所有门店评分明细</div>', unsafe_allow_html=True)
        show_h = health_df.drop(columns=["_color"]).copy()
        st.dataframe(show_h, width='stretch', hide_index=True)


# ── Tab 预警: 预警监控 ────────────────────────────────────────────────────────
with tab_alert:
    alerts_df = compute_alerts(df, z_thresh=z_thresh, drop_thresh=drop_thresh,
                               monitor_metrics=monitor_metrics) if monitor_metrics else pd.DataFrame()
    traffic_alerts = compute_traffic_type_alerts(df)
    refund_alerts  = compute_refund_alerts(df)

    # 顶部总览
    total_metric  = len(alerts_df)
    total_traffic = len(traffic_alerts)
    total_refund  = len(refund_alerts)
    sev_metric    = (alerts_df["严重度"].str.contains("严重").sum()
                     if not alerts_df.empty else 0)
    sev_traffic   = (traffic_alerts["严重度"].str.contains("严重").sum()
                     if not traffic_alerts.empty else 0)
    sev_refund    = (refund_alerts["严重度"].str.contains("严重").sum()
                     if not refund_alerts.empty else 0)

    c1, c2, c3, c4 = st.columns(4)
    with c1: kpi_card("指标异常", str(total_metric), f"严重 {sev_metric} 条")
    with c2: kpi_card("流量类型退化", str(total_traffic), f"严重 {sev_traffic} 条（仅美团）")
    with c3: kpi_card("退单率异常", str(total_refund), f"严重 {sev_refund} 条（仅饿了么）")
    with c4: kpi_card("总预警数", str(total_metric + total_traffic + total_refund), "项")

    sub_a, sub_b, sub_c = st.tabs(["📊 指标异常", "📉 流量类型退化", "↩️ 退单率异常"])

    # ── 子标签 A：指标异常（Z-score + 环比下跌）
    with sub_a:
        st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:14px 18px;
font-size:13px;color:#475569;margin-bottom:8px;'>
<b style='color:#0f172a;'>规则：</b>
Z-score 检测（偏离历史均值 ≥ <b>{z_thresh}σ</b>）+ 环比下跌检测（≥ <b>{int(drop_thresh*100)}%</b>）·
分级：|Z|≥2.5 或下跌≥50% → 🔴 严重 ｜ |Z|≥2.0 或下跌≥35% → 🟠 警告 ｜ 其它 → 🟡 注意
</div>
""", unsafe_allow_html=True)

        if not monitor_metrics:
            st.info("请在左侧选择至少一个监控指标")
        elif alerts_df.empty:
            st.success("✅ 当前筛选范围内未发现明显异常")
        else:
            f1, f2 = st.columns([1, 2])
            with f1:
                sev_filter = st.multiselect(
                    "严重度",
                    ["🔴 严重", "🟠 警告", "🟡 注意"],
                    default=["🔴 严重", "🟠 警告", "🟡 注意"],
                    key="sev_filter_a",
                )
            with f2:
                metric_filter = st.multiselect(
                    "指标",
                    alerts_df["指标"].unique().tolist(),
                    default=alerts_df["指标"].unique().tolist(),
                    key="metric_filter_a",
                )
            show_df = alerts_df[
                alerts_df["严重度"].isin(sev_filter) &
                alerts_df["指标"].isin(metric_filter)
            ][["日期", "门店", "平台", "指标", "当日值", "历史均值", "偏离", "Z值", "严重度"]]
            if show_df.empty:
                st.info("当前筛选下无预警")
            else:
                st.dataframe(show_df, width='stretch', hide_index=True, height=420)
                buf = io.BytesIO()
                with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                    show_df.to_excel(writer, index=False, sheet_name="指标异常")
                buf.seek(0)
                st.download_button(
                    "⬇️ 下载 Excel", data=buf.getvalue(),
                    file_name=f"指标异常_{start_d}_{end_d}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key="dl_alerts_a",
                )

    # ── 子标签 B：流量类型退化（仅美团）
    with sub_b:
        st.markdown("""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:14px 18px;
font-size:13px;color:#475569;margin-bottom:8px;'>
<b style='color:#0f172a;'>规则：</b>美团给每店打的健康度标签层级：
<b>高曝高转 (4)</b> > 高曝低转 (3) > 低曝高转 (2) > 低曝低转 (1)。<br>
检测：当日 vs 7日 退化 / 7日 vs 30日 退化 / 30日类型持续低曝低转 →
跌 1 级 🟠 警告 ｜ 跌 ≥2 级 🔴 严重
</div>
""", unsafe_allow_html=True)
        if traffic_alerts.empty:
            if "美团" not in df["平台"].unique():
                st.info("当前筛选下无美团数据（饿了么无此字段）")
            else:
                st.success("✅ 美团门店未发现流量类型退化")
        else:
            st.dataframe(
                traffic_alerts[["日期", "门店", "平台", "类型", "变化",
                                "当日", "7日", "30日", "严重度"]],
                width='stretch', hide_index=True, height=420,
            )
            buf = io.BytesIO()
            with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                traffic_alerts.drop(columns=["_sev_rank", "_date"], errors="ignore") \
                    .to_excel(writer, index=False, sheet_name="流量类型退化")
            buf.seek(0)
            st.download_button(
                "⬇️ 下载 Excel", data=buf.getvalue(),
                file_name=f"流量类型退化_{start_d}_{end_d}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="dl_alerts_b",
            )

    # ── 子标签 C：退单率异常（仅饿了么）
    with sub_c:
        st.markdown("""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:14px 18px;
font-size:13px;color:#475569;margin-bottom:8px;'>
<b style='color:#0f172a;'>规则：</b>
商户责任退单率 ≥ <b>5%</b> 或 总退单率 ≥ <b>10%</b>（且无效订单 ≥ 3 单）触发预警。<br>
超阈值 1×→ 🟡 注意 ｜ ≥1.5× → 🟠 警告 ｜ ≥2× → 🔴 严重
</div>
""", unsafe_allow_html=True)
        if refund_alerts.empty:
            if "饿了么" not in df["平台"].unique():
                st.info("当前筛选下无饿了么数据（美团表无退单字段）")
            else:
                st.success("✅ 饿了么门店退单率均在正常范围内")
        else:
            st.dataframe(
                refund_alerts[["日期", "门店", "平台", "类型",
                               "退单率", "阈值", "无效订单数", "商户原因数", "严重度"]],
                width='stretch', hide_index=True, height=420,
            )
            buf = io.BytesIO()
            with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                refund_alerts.drop(columns=["_sev_rank", "_date", "_value"], errors="ignore") \
                    .to_excel(writer, index=False, sheet_name="退单率异常")
            buf.seek(0)
            st.download_button(
                "⬇️ 下载 Excel", data=buf.getvalue(),
                file_name=f"退单率异常_{start_d}_{end_d}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="dl_alerts_c",
            )


# ── Tab 利润分析 ──────────────────────────────────────────────────────────────
with tab_profit:
    profit_df = compute_profit_stats(df)
    if profit_df.empty:
        st.info("当前筛选下无数据")
    else:
        # 顶部 KPI：每个平台的总实收/补贴/补贴ROI
        st.markdown('<div class="section-title">利润总览（按平台）</div>', unsafe_allow_html=True)
        plat_kpi_cols = st.columns(len(profit_df["平台"].unique()))
        for i, plat in enumerate(profit_df["平台"].unique()):
            sub = profit_df[profit_df["平台"] == plat]
            t_rev   = sub["实收"].sum()
            t_paid  = sub["客户实付"].sum()
            t_sub   = sub["商家补贴"].sum()
            t_ord   = sub["订单"].sum()
            avg_rev = t_rev / t_ord if t_ord else 0
            roi     = t_rev / t_sub if t_sub else 0
            sub_pct = t_sub / t_paid if t_paid else 0
            color = PLATFORM_COLORS.get(plat, "#888")
            with plat_kpi_cols[i]:
                st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-left:4px solid {color};
border-radius:10px;padding:14px 18px;margin-bottom:10px;'>
<div style='font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px;'>🏷 {plat}</div>
<table style='width:100%;font-size:13px;'>
<tr><td style='color:#64748b;padding:3px 0;'>总实收</td>
    <td style='text-align:right;color:#0f172a;font-weight:700;'>{fmt_money(t_rev)}</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>单均实收</td>
    <td style='text-align:right;color:#0f172a;font-weight:700;'>¥{avg_rev:.2f}</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>商家补贴投入</td>
    <td style='text-align:right;color:#0f172a;font-weight:700;'>{fmt_money(t_sub)}</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>补贴 ROI</td>
    <td style='text-align:right;color:#10b981;font-weight:700;'>{roi:.2f}×</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>补贴占客付比</td>
    <td style='text-align:right;color:#f59e0b;font-weight:700;'>{sub_pct*100:.1f}%</td></tr>
</table></div>
""", unsafe_allow_html=True)

        st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:12px 16px;
font-size:12px;color:#475569;margin-bottom:8px;'>
<b style='color:#0f172a;'>口径说明：</b>
实收 = 平台扣完后到商家账户的钱 · 补贴 ROI = 实收 / 商家自付补贴（越高越赚） ·
补贴占比 = 补贴 / 客户实付（越高补得越狠） · 优惠力度 = 优惠总额 / 优惠前总额
</div>
""", unsafe_allow_html=True)

        # 补贴 ROI 排名图
        st.markdown('<div class="section-title">补贴投入产出比（ROI）排行</div>', unsafe_allow_html=True)
        roi_df = profit_df[profit_df["商家补贴"] > 0].copy()
        if roi_df.empty:
            st.info("当前数据无商家补贴记录")
        else:
            roi_df["门店短名"] = roi_df["门店名称"].apply(_short_store_name)
            roi_df = roi_df.sort_values("补贴ROI", ascending=True)
            fig = go.Figure()
            for plat in roi_df["平台"].unique():
                sub = roi_df[roi_df["平台"] == plat]
                fig.add_trace(go.Bar(
                    y=sub["门店短名"] + f" [{plat}]",
                    x=sub["补贴ROI"],
                    orientation="h",
                    marker_color=PLATFORM_COLORS.get(plat, "#888"),
                    text=[f"{v:.2f}×" for v in sub["补贴ROI"]],
                    textposition="outside",
                    textfont=dict(color=TEXT_COLOR, size=11),
                    name=plat,
                    customdata=sub[["单均补贴", "补贴占比"]],
                    hovertemplate=(
                        "<b>%{y}</b><br>"
                        "补贴 ROI: %{x:.2f}×<br>"
                        "单均补贴: ¥%{customdata[0]:.2f}<br>"
                        "补贴占客付比: %{customdata[1]:.1%}<extra></extra>"
                    ),
                ))
            _base_layout(fig, height=max(380, len(roi_df) * 28))
            fig.update_layout(margin=dict(l=10, r=80, t=40, b=10), showlegend=False, barmode="group")
            fig.update_xaxes(gridcolor=GRID_COLOR, ticksuffix="×", showline=True, linecolor=GRID_COLOR)
            fig.update_yaxes(gridcolor=GRID_COLOR, automargin=True)
            st.plotly_chart(fig, width='stretch')

        # 单均实收 vs 单均补贴 散点图
        st.markdown('<div class="section-title">单均实收 vs 单均补贴（散点）</div>', unsafe_allow_html=True)
        scatter_df = profit_df[profit_df["订单"] > 0].copy()
        scatter_df["门店短名"] = scatter_df["门店名称"].apply(_short_store_name)
        fig2 = px.scatter(
            scatter_df,
            x="单均补贴", y="单均实收",
            color="平台", size="订单",
            color_discrete_map=PLATFORM_COLORS,
            hover_name="门店短名",
            labels={"单均补贴": "单均补贴 (¥)", "单均实收": "单均实收 (¥)"},
        )
        fig2.update_traces(marker=dict(line=dict(color="white", width=1)))
        _base_layout(fig2, height=420)
        fig2.update_xaxes(gridcolor=GRID_COLOR, tickformat=",.1f", showline=True, linecolor=GRID_COLOR, ticksuffix="")
        fig2.update_yaxes(gridcolor=GRID_COLOR, tickformat=",.1f", showline=True, linecolor=GRID_COLOR)
        st.plotly_chart(fig2, width='stretch')

        # 明细表
        st.markdown('<div class="section-title">利润明细表</div>', unsafe_allow_html=True)
        show_p = profit_df.copy()
        show_p = show_p.rename(columns={"门店名称": "门店"})
        show_p["实收"]      = show_p["实收"].map(lambda x: f"¥{x:,.2f}")
        show_p["客户实付"]  = show_p["客户实付"].map(lambda x: f"¥{x:,.2f}")
        show_p["商家补贴"]  = show_p["商家补贴"].map(lambda x: f"¥{x:,.2f}")
        show_p["平台费"]    = show_p["平台费"].map(lambda x: f"¥{x:,.2f}")
        show_p["单均实收"]  = show_p["单均实收"].map(lambda x: f"¥{x:.2f}")
        show_p["单均补贴"]  = show_p["单均补贴"].map(lambda x: f"¥{x:.2f}")
        show_p["实收率"]    = show_p["实收率"].map(lambda x: f"{x*100:.1f}%")
        show_p["补贴ROI"]   = show_p["补贴ROI"].map(lambda x: f"{x:.2f}×" if x else "—")
        show_p["补贴占比"]  = show_p["补贴占比"].map(lambda x: f"{x*100:.1f}%")
        show_p["优惠力度"]  = show_p["优惠力度"].map(lambda x: f"{x*100:.1f}%")
        show_p["订单"]      = show_p["订单"].map(lambda x: f"{int(x):,}")
        show_p = show_p[["门店", "平台", "订单", "实收", "单均实收",
                         "商家补贴", "单均补贴", "补贴ROI", "补贴占比",
                         "实收率", "优惠力度", "平台费"]]
        st.dataframe(show_p, width='stretch', hide_index=True)


# ── Tab 大盘对比: 门店 vs 平台基准 ────────────────────────────────────────────
with tab_peer:
    benchmarks, peer_df = compute_peer_comparison(df)

    if not benchmarks:
        st.info("当前筛选下无数据")
    else:
        # 顶部：各平台大盘基准卡片
        st.markdown('<div class="section-title">各平台大盘基准</div>', unsafe_allow_html=True)
        plat_list = list(benchmarks.keys())
        bench_cols = st.columns(len(plat_list))
        for i, plat in enumerate(plat_list):
            b = benchmarks[plat]
            color = PLATFORM_COLORS.get(plat, "#888")
            with bench_cols[i]:
                st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-left:4px solid {color};
border-radius:10px;padding:14px 18px;margin-bottom:10px;'>
<div style='font-size:14px;font-weight:700;color:#0f172a;margin-bottom:10px;'>🏷 {plat} 大盘</div>
<table style='width:100%;font-size:13px;'>
<tr><td style='color:#64748b;padding:3px 0;'>下单转化率</td>
<td style='text-align:right;color:#0f172a;font-weight:700;'>{b['下单转化率']*100:.2f}%</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>进店转化率</td>
<td style='text-align:right;color:#0f172a;font-weight:700;'>{b['进店转化率']*100:.2f}%</td></tr>
<tr><td style='color:#64748b;padding:3px 0;'>客单价</td>
<td style='text-align:right;color:#0f172a;font-weight:700;'>¥{b['客单价']:.2f}</td></tr>
<tr><td style='color:#94a3b8;padding:3px 0;font-size:11px;'>样本</td>
<td style='text-align:right;color:#94a3b8;font-size:11px;'>{int(b['_order_total']):,} 单 / ¥{b['_rev_total']:,.0f}</td></tr>
</table>
</div>
""", unsafe_allow_html=True)

        st.markdown(f"""
<div style='background:#fff;border:1px solid #e2e8f0;border-radius:8px;padding:12px 16px;
font-size:12px;color:#475569;margin-bottom:8px;'>
<b style='color:#0f172a;'>计算方式：</b>
大盘=该平台所有门店加权平均 ·
下单转化率=总订单/总进店 · 进店转化率=总进店/总曝光 · 客单价=总收入/总订单。
偏离≤-30% 严重偏低 · -30%~-15% 偏低 · ±15%内 接近大盘 · ≥+15% 表现良好 · ≥+30% 标杆。
</div>
""", unsafe_allow_html=True)

        # 指标选择
        col_m, col_p = st.columns([2, 1])
        with col_m:
            sel_metric = st.radio(
                "对比指标",
                ["下单转化率", "进店转化率", "客单价"],
                horizontal=True,
                key="peer_metric",
            )
        with col_p:
            sel_plat = st.selectbox(
                "平台",
                ["全部"] + plat_list,
                key="peer_plat",
            )

        plot_df = peer_df[peer_df["指标"] == sel_metric].copy()
        if sel_plat != "全部":
            plot_df = plot_df[plot_df["平台"] == sel_plat]

        if plot_df.empty:
            st.info("无数据")
        else:
            plot_df["分级"] = plot_df["偏离"].apply(lambda x: classify_peer_deviation(x)[0])
            plot_df["颜色"] = plot_df["偏离"].apply(lambda x: classify_peer_deviation(x)[1])
            plot_df["门店短名"] = plot_df["门店"].apply(_short_store_name)
            plot_df = plot_df.sort_values("偏离", ascending=True)

            # 图表：条形图，参考线 = 大盘均值
            is_pct = plot_df["_is_pct"].iloc[0]

            fig = go.Figure()
            for plat in plot_df["平台"].unique():
                sub = plot_df[plot_df["平台"] == plat]
                bench = benchmarks[plat][sel_metric]
                if is_pct:
                    x_vals = sub["门店值"] * 100
                    bench_v = bench * 100
                    suffix = "%"
                else:
                    x_vals = sub["门店值"]
                    bench_v = bench
                    suffix = ""
                fig.add_trace(go.Bar(
                    y=sub["门店短名"] + f" [{plat}]",
                    x=x_vals,
                    orientation="h",
                    marker_color=sub["颜色"],
                    text=[f"{v:.2f}{suffix}" for v in x_vals],
                    textposition="outside",
                    textfont=dict(color=TEXT_COLOR, size=11),
                    name=plat,
                    customdata=sub[["偏离", "分级"]],
                    hovertemplate=(
                        "<b>%{y}</b><br>"
                        f"{sel_metric}: %{{x:.2f}}{suffix}<br>"
                        "偏离大盘: %{customdata[0]:+.1%}<br>"
                        "分级: %{customdata[1]}<extra></extra>"
                    ),
                ))
                # 添加大盘参考线
                fig.add_vline(
                    x=bench_v,
                    line=dict(color=PLATFORM_COLORS.get(plat, "#888"),
                              width=2, dash="dash"),
                    annotation_text=f"{plat} 大盘 {bench_v:.2f}{suffix}",
                    annotation_position="top",
                    annotation_font_color=PLATFORM_COLORS.get(plat, "#888"),
                )

            _base_layout(fig, height=max(380, len(plot_df) * 28))
            fig.update_layout(margin=dict(l=10, r=80, t=60, b=10), showlegend=False)
            fig.update_xaxes(gridcolor=GRID_COLOR, ticksuffix=suffix, showline=True, linecolor=GRID_COLOR)
            fig.update_yaxes(gridcolor=GRID_COLOR, automargin=True)

            st.markdown(f'<div class="section-title">{sel_metric} · 各门店 vs 大盘</div>', unsafe_allow_html=True)
            st.plotly_chart(fig, width='stretch')

            # 偏低门店（重点关注）
            problem = plot_df[plot_df["偏离"] <= -0.15].sort_values("偏离")
            top = plot_df[plot_df["偏离"] >= 0.15].sort_values("偏离", ascending=False)

            col_l, col_r = st.columns(2)
            with col_l:
                st.markdown('<div class="section-title">⚠️ 偏低门店（需关注）</div>', unsafe_allow_html=True)
                if problem.empty:
                    st.success("✅ 所有门店在该指标上均接近或高于大盘")
                else:
                    show = problem[["门店", "平台", "门店值", "大盘均值", "偏离", "分级"]].copy()
                    if is_pct:
                        show["门店值"]  = show["门店值"].map(lambda x: f"{x*100:.2f}%")
                        show["大盘均值"] = show["大盘均值"].map(lambda x: f"{x*100:.2f}%")
                    else:
                        show["门店值"]  = show["门店值"].map(lambda x: f"¥{x:.2f}")
                        show["大盘均值"] = show["大盘均值"].map(lambda x: f"¥{x:.2f}")
                    show["偏离"] = show["偏离"].map(lambda x: f"{x*100:+.1f}%")
                    st.dataframe(show, width='stretch', hide_index=True)

            with col_r:
                st.markdown('<div class="section-title">🏆 标杆门店（可借鉴）</div>', unsafe_allow_html=True)
                if top.empty:
                    st.info("当前指标无明显标杆门店")
                else:
                    show = top[["门店", "平台", "门店值", "大盘均值", "偏离", "分级"]].copy()
                    if is_pct:
                        show["门店值"]  = show["门店值"].map(lambda x: f"{x*100:.2f}%")
                        show["大盘均值"] = show["大盘均值"].map(lambda x: f"{x*100:.2f}%")
                    else:
                        show["门店值"]  = show["门店值"].map(lambda x: f"¥{x:.2f}")
                        show["大盘均值"] = show["大盘均值"].map(lambda x: f"¥{x:.2f}")
                    show["偏离"] = show["偏离"].map(lambda x: f"{x*100:+.1f}%")
                    st.dataframe(show, width='stretch', hide_index=True)


# ── Tab 4: 日报 ───────────────────────────────────────────────────────────────
with tab4:
    st.markdown('<div class="section-title">选择日期生成日报</div>', unsafe_allow_html=True)

    report_date = st.date_input(
        "报告日期",
        value=max_d,
        min_value=min_d,
        max_value=max_d,
        key="report_date_input",
    )

    if isinstance(report_date, (tuple, list)):
        report_date = report_date[0]

    report_df = generate_daily_report(df_all, report_date)

    if report_df.empty:
        st.warning(f"⚠️ {report_date} 无数据")
    else:
        day_data = df_all[df_all["日期"] == pd.Timestamp(report_date)]
        total_rev_day = day_data["收入"].sum()
        total_ord_day = day_data["有效订单"].sum()
        total_exp_day = day_data["曝光人数"].sum() if "曝光人数" in day_data.columns else 0
        store_cnt_day = day_data["门店名称"].nunique()
        avg_price_day = (total_rev_day / total_ord_day) if total_ord_day else 0

        st.markdown(f"""
<div class="report-header">
<h3>📋 {report_date} 经营日报</h3>
<p style='margin:0;font-size:14px;'>
  营业收入 <b>{fmt_money(total_rev_day)}</b> &nbsp;|&nbsp;
  有效订单 <b>{fmt_int(total_ord_day)} 单</b> &nbsp;|&nbsp;
  客单价 <b>¥{avg_price_day:.2f}</b> &nbsp;|&nbsp;
  曝光人数 <b>{fmt_int(total_exp_day)}</b> &nbsp;|&nbsp;
  活跃门店 <b>{store_cnt_day} 家</b>
</p>
</div>
""", unsafe_allow_html=True)

        st.dataframe(report_df, width='stretch', hide_index=True)

        # 下载按钮组
        st.markdown('<div class="section-title">📥 导出日报</div>', unsafe_allow_html=True)
        dl_c1, dl_c2 = st.columns(2)

        with dl_c1:
            # Excel
            excel_buf = io.BytesIO()
            with pd.ExcelWriter(excel_buf, engine="openpyxl") as writer:
                report_df.to_excel(writer, index=False, sheet_name=str(report_date))
            excel_buf.seek(0)
            st.download_button(
                label="📊 下载 Excel 日报",
                data=excel_buf.getvalue(),
                file_name=f"日报_{report_date}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key=f"dl_excel_{report_date}",
                width='stretch',
            )

        with dl_c2:
            # PPT - 实时生成（含当日预警）
            try:
                ppt_alerts = compute_alerts(df_all, z_thresh=z_thresh,
                                             drop_thresh=drop_thresh,
                                             monitor_metrics=monitor_metrics) if monitor_metrics else pd.DataFrame()
                ppt_traffic = compute_traffic_type_alerts(df_all)
                ppt_refund  = compute_refund_alerts(df_all)
                ppt_buf = generate_ppt_report(df_all, report_date,
                                                ppt_alerts, ppt_traffic, ppt_refund)
                st.download_button(
                    label="🎬 下载 PPT 日报",
                    data=ppt_buf.getvalue(),
                    file_name=f"日报_{report_date}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"dl_ppt_{report_date}",
                    width='stretch',
                )
            except Exception as e:
                st.error(f"PPT 生成失败：{e}")

        st.caption("💡 PPT 日报包含 5 张幻灯片：封面、核心指标、Top 5 门店、需关注门店、预警汇总。可直接给老板看。")


# ── Tab 到手率分析 ────────────────────────────────────────────────────────────
with tab_rates:
    st.markdown('<div class="section-title">数据文件配置</div>', unsafe_allow_html=True)

    # 支持通过环境变量设置默认数据文件夹，实现"放进去自动读"
    _RATES_DEFAULT_DIR = os.environ.get("RATES_DATA_DIR", "")

    with st.expander("📁 文件配置", expanded=True):
        ra_col1, ra_col2 = st.columns([2, 1])
        with ra_col1:
            ra_dir = st.text_input(
                "数据文件夹路径（配置后自动读取，留空则手动上传）",
                value=_RATES_DEFAULT_DIR, key="ra_dir",
                help="把 4 个 Excel 文件放到同一文件夹，填入路径后刷新页面即可自动加载。",
            )
        with ra_col2:
            ra_month = st.text_input("月份标签（用于报告标题）", value="2026-04", key="ra_month")
            ra_delivery = st.number_input(
                "小程序配送费（元）", min_value=0.0, value=0.0, step=0.01, key="ra_delivery",
                help="当月小程序外卖配送费总额，在汇总层面整体扣减。",
            )

        up_c1, up_c2, up_c3, up_c4 = st.columns(4)
        with up_c1:
            ra_main_up = st.file_uploader("主表（含门店价）", type=["xlsx", "xls"], key="ra_main_up")
            ra_main_name = st.text_input("文件名", value="主表.xlsx", key="ra_main_name", label_visibility="collapsed")
        with up_c2:
            ra_mt_up = st.file_uploader("美团渠道表", type=["xlsx", "xls"], key="ra_mt_up")
            ra_mt_name = st.text_input("文件名", value="美团.xlsx", key="ra_mt_name", label_visibility="collapsed")
        with up_c3:
            ra_ele_up = st.file_uploader("饿了么渠道表", type=["xlsx", "xls"], key="ra_ele_up")
            ra_ele_name = st.text_input("文件名", value="饿了么.xlsx", key="ra_ele_name", label_visibility="collapsed")
        with up_c4:
            ra_mp_up = st.file_uploader("小程序订单流水", type=["xlsx", "xls"], key="ra_mp_up")
            ra_mp_name = st.text_input("文件名", value="小程序.xlsx", key="ra_mp_name", label_visibility="collapsed")

    def _ra_src(upload, folder, name):
        if upload is not None:
            return upload
        if folder:
            p = os.path.join(folder, name)
            if os.path.exists(p):
                return p
        return None

    ra_main_src = _ra_src(ra_main_up, ra_dir, ra_main_name)
    ra_mt_src   = _ra_src(ra_mt_up,   ra_dir, ra_mt_name)
    ra_ele_src  = _ra_src(ra_ele_up,  ra_dir, ra_ele_name)
    ra_mp_src   = _ra_src(ra_mp_up,   ra_dir, ra_mp_name)

    ra_missing = [n for n, s in [("主表", ra_main_src), ("美团", ra_mt_src),
                                  ("饿了么", ra_ele_src), ("小程序", ra_mp_src)] if s is None]
    if ra_missing:
        st.info(f"👆 请上传或配置数据文件夹路径，缺少：{', '.join(ra_missing)}")
        st.caption("提示：把 4 个 Excel 文件放到同一文件夹，在上方填入路径，每次更新文件后刷新页面即自动加载。")
    else:
        def _fp(v):
            return f"{v * 100:.2f}%" if v is not None else "-"

        def _fm(v):
            return f"¥{v:,.2f}" if v is not None else "-"

        with st.spinner("正在读取并计算..."):
            ra_msgs, ra_errs = [], []
            ra_ok = False
            try:
                pdict, has_spec, main_n = rates_load_main(ra_main_src)
                ra_msgs.append(("✅", f"主表读取成功，共 {main_n} 个 SKU{'（含规格）' if has_spec else '（无规格列）'}"))

                mt_df, mt_s = rates_load_channel(ra_mt_src, pdict, has_spec, "美团")
                n_mt, nm_mt = mt_s['sku_count'], mt_s['matched_count']
                icon_mt = "✅" if not mt_s['unmatched'] else "⚠️"
                ra_msgs.append((icon_mt, f"美团渠道：{n_mt} 个成品 SKU，门店价匹配率 {nm_mt/n_mt*100:.0f}%" if n_mt else "美团：无成品 SKU"))

                ele_df, ele_s = rates_load_channel(ra_ele_src, pdict, has_spec, "饿了么")
                n_el, nm_el = ele_s['sku_count'], ele_s['matched_count']
                icon_el = "✅" if not ele_s['unmatched'] else "⚠️"
                ra_msgs.append((icon_el, f"饿了么渠道：{n_el} 个成品 SKU，门店价匹配率 {nm_el/n_el*100:.0f}%" if n_el else "饿了么：无成品 SKU"))

                mp_df, mp_s = rates_load_mp(ra_mp_src, float(ra_delivery))
                ra_msgs.append(("✅", f"小程序：{mp_s['order_count']} 笔外卖订单，已剔除 {mp_s['removed']} 行汇总数据"))
                if mp_s.get('warning'):
                    ra_msgs.append(("⚠️", mp_s['warning']))

                t_orig = mt_s['系统销售原价'] + ele_s['系统销售原价'] + mp_s['系统销售原价']
                t_real = mt_s['门店真实销售额'] + ele_s['门店真实销售额'] + mp_s['门店真实销售额']
                t_recv = mt_s['实收金额'] + ele_s['实收金额'] + mp_s['净实收']
                ov = {
                    '系统销售原价': t_orig, '门店真实销售额': t_real, '实收金额': t_recv,
                    '系统到手率': t_recv / t_orig if t_orig > 0 else None,
                    '真实到手率': t_recv / t_real if t_real > 0 else None,
                    '溢价率': (t_orig - t_real) / t_real if t_real > 0 else None,
                }
                ra_ok = True
            except Exception as e:
                ra_errs.append(str(e))

        for icon, msg in ra_msgs:
            if icon == "✅":
                st.success(msg)
            else:
                st.warning(msg)
        for err in ra_errs:
            st.error(err)

        if ra_ok:
            # ── 概览卡片 ──────────────────────────────────────────────────────
            st.markdown(f'<div class="section-title">{ra_month} 到手率概览</div>', unsafe_allow_html=True)
            rc1, rc2, rc3, rc4 = st.columns(4)

            def _ra_card(col, label, main, subs):
                col.markdown(
                    f"""<div class="metric-card">
                    <div class="metric-label">{label}</div>
                    <div class="metric-value">{main}</div>
                    {''.join(f'<div class="metric-sub">{s}</div>' for s in subs)}
                    </div>""",
                    unsafe_allow_html=True,
                )

            _ra_card(rc1, f"{ra_month} 整体 真实到手率", _fp(ov['真实到手率']),
                     [f"系统到手率: {_fp(ov['系统到手率'])}", f"溢价率: {_fp(ov['溢价率'])}", f"实收: {_fm(ov['实收金额'])}"])
            _ra_card(rc2, "美团 真实到手率", _fp(mt_s['真实到手率']),
                     [f"系统到手率: {_fp(mt_s['系统到手率'])}", f"溢价率: {_fp(mt_s['溢价率'])}", f"成品SKU: {mt_s['sku_count']}"])
            _ra_card(rc3, "饿了么 真实到手率", _fp(ele_s['真实到手率']),
                     [f"系统到手率: {_fp(ele_s['系统到手率'])}", f"溢价率: {_fp(ele_s['溢价率'])}", f"成品SKU: {ele_s['sku_count']}"])
            _ra_card(rc4, "小程序 净到手率（扣配送）", _fp(mp_s['净到手率']),
                     [f"毛到手率: {_fp(mp_s['毛到手率'])}", f"配送费: {_fm(mp_s['配送费'])}", f"订单数: {mp_s['order_count']}"])

            # ── 分渠道汇总表 ──────────────────────────────────────────────────
            st.markdown('<div class="section-title">分渠道汇总</div>', unsafe_allow_html=True)
            ch_rows = []
            for ch_name, s, qty_label in [
                ("美团",   mt_s,  f"{mt_s['销售数量']:.0f} 件"),
                ("饿了么", ele_s, f"{ele_s['销售数量']:.0f} 件"),
                ("小程序", mp_s,  f"{mp_s['order_count']} 单"),
            ]:
                recv_disp = (_fm(s['净实收']) + " (扣配送)") if ch_name == "小程序" else _fm(s['实收金额'])
                ch_rows.append({
                    "渠道": ch_name, "数量/订单": qty_label,
                    "系统销售原价": _fm(s['系统销售原价']),
                    "门店真实销售额": _fm(s['门店真实销售额']),
                    "实收金额": recv_disp,
                    "折扣金额": _fm(s['折扣金额']),
                    "系统到手率": _fp(s['系统到手率']),
                    "真实到手率": _fp(s['真实到手率']),
                    "溢价率": _fp(s['溢价率']),
                })
            ch_rows.append({
                "渠道": "合计", "数量/订单": "-",
                "系统销售原价": _fm(ov['系统销售原价']),
                "门店真实销售额": _fm(ov['门店真实销售额']),
                "实收金额": _fm(ov['实收金额']),
                "折扣金额": "-",
                "系统到手率": _fp(ov['系统到手率']),
                "真实到手率": _fp(ov['真实到手率']),
                "溢价率": _fp(ov['溢价率']),
            })
            st.dataframe(pd.DataFrame(ch_rows), hide_index=True, use_container_width=True)

            # ── SKU / 订单明细 ────────────────────────────────────────────────
            st.markdown('<div class="section-title">明细数据</div>', unsafe_allow_html=True)
            rt1, rt2, rt3 = st.tabs(["🛵 美团 SKU", "🟡 饿了么 SKU", "📱 小程序订单"])

            def _show_ch_detail(tab, df, summ):
                with tab:
                    if summ['unmatched']:
                        names = "；".join(f"{u['末级分类']}/{u['商品名称']}" for u in summ['unmatched'])
                        st.warning(f"以下 SKU 未匹配到门店价（将影响真实到手率计算）：{names}")
                    disp = df.copy()
                    for col in ['系统到手率', '真实到手率', '溢价率']:
                        if col in disp.columns:
                            disp[col] = disp[col].apply(lambda v: _fp(v) if pd.notna(v) else "-")
                    st.dataframe(disp, hide_index=True, use_container_width=True, height=420)

            _show_ch_detail(rt1, mt_df, mt_s)
            _show_ch_detail(rt2, ele_df, ele_s)

            with rt3:
                st.caption(
                    f"毛到手率: **{_fp(mp_s['毛到手率'])}**  |  "
                    f"配送费: **{_fm(mp_s['配送费'])}**  |  "
                    f"净实收: **{_fm(mp_s['净实收'])}**  |  "
                    f"净到手率: **{_fp(mp_s['净到手率'])}**"
                )
                st.dataframe(mp_df, hide_index=True, use_container_width=True, height=420)

            # ── 导出 Excel ────────────────────────────────────────────────────
            st.markdown('<div class="section-title">📥 导出 Excel 报告</div>', unsafe_allow_html=True)
            ra_excel_buf = io.BytesIO()
            with pd.ExcelWriter(ra_excel_buf, engine='openpyxl') as writer:
                pd.DataFrame(ch_rows).to_excel(writer, index=False, sheet_name="渠道汇总")
                mt_df.to_excel(writer, index=False, sheet_name=f"{ra_month}美团")
                ele_df.to_excel(writer, index=False, sheet_name=f"{ra_month}饿了么")
                mp_df.to_excel(writer, index=False, sheet_name=f"{ra_month}小程序(外卖)")
            ra_excel_buf.seek(0)
            st.download_button(
                label="📊 下载 Excel 报告",
                data=ra_excel_buf.getvalue(),
                file_name=f"到手率分析_{ra_month}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="ra_dl_excel",
            )
